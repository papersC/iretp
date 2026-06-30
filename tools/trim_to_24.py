"""Remove slides 25-27 (the appended Architecture/CLA/Hardware slides) so we can rebuild."""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pptx import Presentation
from pptx.oxml.ns import qn

SRC = r"C:\Users\kalmi\IRETP\IRETP_Technical_Presentation.pptx"

def delete_slide(prs, index):
    """index is 0-based"""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    target = slides[index]
    # Remove from rels
    rId = target.get(qn("r:id"))
    prs.part.drop_rel(rId)
    # Remove from sldIdLst
    xml_slides.remove(target)

prs = Presentation(SRC)
count = len(prs.slides)
print(f"Before: {count} slides")
# Remove last (count-24) slides
while len(prs.slides) > 24:
    delete_slide(prs, len(prs.slides) - 1)
print(f"After:  {len(prs.slides)} slides")
prs.save(SRC)
