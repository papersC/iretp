"""Build IRETP_Technical_Presentation.pptx — a verified technical deck that
mirrors the four reports under reports/ and reflects the actual state of the
codebase under src/. All facts are cross-checked against the source tree.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from lxml import etree

# --- Palette (DLD-green inspired forest theme) -----------------------------
FOREST       = RGBColor(0x0E, 0x4E, 0x37)  # primary dark green
FOREST_DEEP  = RGBColor(0x06, 0x35, 0x25)  # title backdrops
SAGE         = RGBColor(0x7C, 0xB9, 0x9F)  # secondary
SAGE_SOFT    = RGBColor(0xE6, 0xF1, 0xEB)  # tint panel
GOLD         = RGBColor(0xC8, 0xA4, 0x64)  # accent (Dubai sand)
GOLD_SOFT    = RGBColor(0xF3, 0xE9, 0xD2)
INK          = RGBColor(0x14, 0x1A, 0x18)
INK_SOFT     = RGBColor(0x4A, 0x55, 0x52)
MUTED        = RGBColor(0x8B, 0x94, 0x91)
PAPER        = RGBColor(0xFF, 0xFF, 0xFF)
DIVIDER      = RGBColor(0xD7, 0xDE, 0xDA)

# Fonts
HEAD_FONT = "Georgia"
BODY_FONT = "Calibri"

OUT          = Path(r"C:\Users\kalmi\IRETP\IRETP_Technical_Presentation.pptx")
OUT_INTERNAL = Path(r"C:\Users\kalmi\IRETP\IRETP_Internal_Costs.pptx")


# --- Helpers ---------------------------------------------------------------
def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def add_rect(slide, x, y, w, h, fill, line=None, line_width=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_width:
            shp.line.width = line_width
    shp.shadow.inherit = False
    return shp


def add_round(slide, x, y, w, h, fill, corner=0.06, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = corner
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, italic=False,
             color=INK, font=BODY_FONT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = ln
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=13, color=INK,
                bullet_color=FOREST, line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # leading bullet (square) in accent color, then space, then body
        r_bullet = p.add_run()
        r_bullet.text = "■  "
        r_bullet.font.name = BODY_FONT
        r_bullet.font.size = Pt(size)
        r_bullet.font.bold = True
        r_bullet.font.color.rgb = bullet_color
        r_body = p.add_run()
        r_body.text = item
        r_body.font.name = BODY_FONT
        r_body.font.size = Pt(size)
        r_body.font.color.rgb = color
    return tb


def add_kicker(slide, x, y, w, text, color=GOLD):
    add_text(slide, x, y, w, Inches(0.3), text.upper(),
             size=10, bold=True, color=color, font=BODY_FONT, line_spacing=1.0)


def add_title(slide, x, y, w, h, text, *, color=FOREST):
    add_text(slide, x, y, w, h, text,
             size=30, bold=True, color=color, font=HEAD_FONT, line_spacing=1.05)


def add_footer(slide, prs, num=None, total=None, label=None):
    sw = prs.slide_width
    sh = prs.slide_height
    add_text(slide, Inches(0.5), sh - Inches(0.36), Inches(8.0), Inches(0.25),
             label or "IRETP — Technical Presentation",
             size=9, color=MUTED, font=BODY_FONT, line_spacing=1.0)
    if num is not None and total is not None:
        add_text(slide, sw - Inches(1.3), sh - Inches(0.36),
                 Inches(0.9), Inches(0.25),
                 f"{num} / {total}",
                 size=9, color=MUTED, font=BODY_FONT,
                 align=PP_ALIGN.RIGHT, line_spacing=1.0)


def add_header(slide, prs, kicker, title, *, kicker_color=GOLD):
    """Standard page header: gold kicker + green serif title."""
    add_kicker(slide, Inches(0.5), Inches(0.45), Inches(12), kicker, color=kicker_color)
    add_title(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.7), title)
    # Subtle baseline rule under the title area (full-width is fine on a
    # single thin line; it visually anchors the page, isn't a colored bar.)
    add_rect(slide, Inches(0.5), Inches(1.5), Inches(12.3), Emu(9525), DIVIDER)


def add_table(slide, x, y, w, h, headers, rows, *,
              header_fill=FOREST, header_fg=PAPER, body_fg=INK,
              row_alt=SAGE_SOFT, first_col_bold=True, font_size=11,
              col_widths=None):
    cols = len(headers)
    rs = len(rows) + 1
    tbl_shape = slide.shapes.add_table(rs, cols, x, y, w, h)
    tbl = tbl_shape.table
    if col_widths is not None:
        assert len(col_widths) == cols
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = cw
    # headers
    for ci, hd in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.margin_left = Inches(0.1)
        cell.margin_right = Inches(0.1)
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)
        tf = cell.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = hd
        run.font.name = BODY_FONT
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = header_fg
    # rows
    for ri, row in enumerate(rows, start=1):
        for ci, cellv in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (row_alt if ri % 2 == 1 else PAPER)
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.1
            run = p.add_run()
            run.text = str(cellv)
            run.font.name = BODY_FONT
            run.font.size = Pt(font_size)
            run.font.bold = first_col_bold and ci == 0
            run.font.color.rgb = body_fg
    return tbl_shape


def add_left_accent(slide, top=None, height=None, color=GOLD):
    """A thin vertical accent rule on the left margin — our visual motif."""
    h = height or slide.part.package.presentation_part.presentation.slide_height
    t = top or 0
    add_rect(slide, Inches(0.18), t, Emu(38100), h, color)


# --- Slides ----------------------------------------------------------------
def slide_cover(prs):
    s = add_blank(prs)
    sw = prs.slide_width
    sh = prs.slide_height
    add_rect(s, 0, 0, sw, sh, FOREST_DEEP)
    # Soft sage band — top right
    add_round(s, sw - Inches(3.8), -Inches(0.8), Inches(4.2), Inches(3.0),
              FOREST, corner=0.5)
    # Gold accent rail on the left
    add_rect(s, Inches(0.0), Inches(0.0), Emu(57150), sh, GOLD)

    add_text(s, Inches(0.9), Inches(0.95), Inches(9), Inches(0.35),
             "DUBAI LAND DEPARTMENT  ·  RFP DLD-IRETP-2026-001",
             size=11, bold=True, color=GOLD, font=BODY_FONT, line_spacing=1.0)
    add_text(s, Inches(0.9), Inches(1.4), Inches(11), Inches(2.5),
             "Integrated Real Estate\nTransparency Platform",
             size=54, bold=True, color=PAPER, font=HEAD_FONT, line_spacing=1.05)
    add_text(s, Inches(0.9), Inches(4.2), Inches(11), Inches(0.6),
             "Technical Presentation — Implementation & Compliance Briefing",
             size=20, italic=True, color=SAGE, font=HEAD_FONT, line_spacing=1.1)

    # Bottom-left meta block
    add_text(s, Inches(0.9), Inches(6.05), Inches(6), Inches(0.3),
             "PROGRAMME",
             size=9, bold=True, color=GOLD, font=BODY_FONT, line_spacing=1.0)
    add_text(s, Inches(0.9), Inches(6.30), Inches(7), Inches(0.4),
             "IRETP  ·  Production-grade Clean-Architecture .NET 9 platform",
             size=13, color=PAPER, font=BODY_FONT, line_spacing=1.1)

    add_text(s, Inches(0.9), Inches(6.75), Inches(7), Inches(0.3),
             "SCOPE OF THIS DECK",
             size=9, bold=True, color=GOLD, font=BODY_FONT, line_spacing=1.0)
    add_text(s, Inches(0.9), Inches(7.0), Inches(11.5), Inches(0.4),
             "Technical Proposal · Compliance Matrix · DESC ISR v3 Appendix · Financial Proposal · Code Verification",
             size=12, color=PAPER, font=BODY_FONT, line_spacing=1.1)


def slide_programme(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "Programme overview", "Why IRETP, and why this proposal stands apart")
    # Two columns
    left_x = Inches(0.5); col_w = Inches(6.1); col_y = Inches(1.85)
    right_x = Inches(6.8)

    add_kicker(s, left_x, col_y, col_w, "What IRETP is")
    add_text(s, left_x, col_y + Inches(0.3), col_w, Inches(0.5),
             "A single transparency platform for Dubai real-estate",
             size=18, bold=True, color=INK, font=HEAD_FONT)
    add_text(s, left_x, col_y + Inches(0.95), col_w, Inches(3.5),
             "IRETP unifies the external investor experience, internal DLD "
             "oversight and the Open Data ecosystem into one Clean-Architecture "
             ".NET 9 solution.\n\n"
             "Three deployable hosts — IRETP.Web (Blazor Server portal), "
             "IRETP.WebAPI (public + investor APIs), and IRETP.AdminAPI "
             "(internal DLD operations) — backed by a shared CQRS Application "
             "layer (MediatR), DDD Domain model and Infrastructure that hosts "
             "EF Core, OneLake, Hangfire, the AI orchestration suite and "
             "identity.",
             size=13, color=INK_SOFT, line_spacing=1.35)

    add_kicker(s, right_x, col_y, col_w, "Why this proposal is different")
    add_text(s, right_x, col_y + Inches(0.3), col_w, Inches(0.5),
             "A live codebase, not a greenfield plan",
             size=18, bold=True, color=INK, font=HEAD_FONT)
    add_bullets(s, right_x, col_y + Inches(1.05), col_w, Inches(4.5), [
        "Every RFP clause maps to a discrete feature folder under src/ and a Razor page or controller already in the repository.",
        "Production-grade engineering posture: 0 build warnings, 0 errors on a clean solution build, structured logging via Serilog.",
        "Strict inward dependency rule: Application has no EF Core; Domain has no Infrastructure. Audited file-by-file.",
        "All AI inference and personal data planned to remain UAE-resident (Azure UAE North or UAE-hosted OSS LLM).",
    ], size=13, color=INK_SOFT)

    add_footer(s, prs, num, total)


def slide_documents(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "Document set", "Four reports + this presentation")

    cards = [
        ("Vol. 1", "Technical Solution Proposal",
         "Clause-by-clause solution narrative across all 20 RFP sections; "
         "file paths and feature folders for every implementation claim; "
         "Section 18 lists 15 engineering commitments included within the delivery scope."),
        ("Vol. 2", "RFP Compliance Matrix",
         "Status of every requirement — C (compliant today), CD (compliant by "
         "delivery, foundations in repo), E (enhanced beyond requirement). "
         "Summary roll-up by RFP section and the value-add commitments group."),
        ("Apx. A", "DESC ISR v3 Control Mapping",
         "35 controls across 10 domains. 27 Implemented, 3 Partially "
         "Implemented (runbook items), 5 Planned (pre-VAPT). Used as the "
         "vendor's pre-VAPT self-assessment for Phase 3 DESC submission."),
        ("Vol. 3", "Financial Proposal",
         "Indicative CAPEX AED 2.58 M (Option A, Azure UAE North) and annual "
         "OPEX AED 629 K. 6 payment milestones tied to phase-gate UAT "
         "sign-offs across 4 months. Three costed hosting options."),
    ]
    # 2x2 grid of cards
    grid_x = Inches(0.5); grid_y = Inches(1.85)
    card_w = Inches(6.1); card_h = Inches(2.45); gap = Inches(0.2)
    for i, (tag, name, body) in enumerate(cards):
        r = i // 2; c = i % 2
        cx = grid_x + (card_w + gap) * c
        cy = grid_y + (card_h + gap) * r
        add_round(s, cx, cy, card_w, card_h, PAPER, corner=0.04, line=DIVIDER)
        # Left tag column
        add_rect(s, cx, cy, Inches(1.05), card_h, FOREST)
        add_text(s, cx, cy + Inches(0.9), Inches(1.05), Inches(0.5),
                 tag, size=18, bold=True, color=GOLD, font=HEAD_FONT,
                 align=PP_ALIGN.CENTER)
        # Right content
        add_text(s, cx + Inches(1.25), cy + Inches(0.18), card_w - Inches(1.4),
                 Inches(0.45), name,
                 size=15, bold=True, color=INK, font=HEAD_FONT)
        add_text(s, cx + Inches(1.25), cy + Inches(0.68),
                 card_w - Inches(1.4), Inches(1.7), body,
                 size=11, color=INK_SOFT, line_spacing=1.35)

    add_footer(s, prs, num, total)


def slide_verification(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "Verification", "Reports cross-checked against src/ today")

    # Top: four big stat tiles
    stats = [
        (".NET 9.0", "Target framework",        "All 6 .csproj files (reports said .NET 8)"),
        ("28",       "API controllers",         "20 in WebAPI · 8 in AdminAPI"),
        ("9",        "Hangfire recurring jobs", "Risk · Score · Alerts · KPI · Escrow · Currency · Benchmark"),
        ("0",        "Build warnings & errors", "Clean solution build"),
    ]
    tile_x = Inches(0.5); tile_y = Inches(1.85)
    tile_w = Inches(3.07); tile_h = Inches(1.65); gap = Inches(0.13)
    for i, (big, lbl, sub) in enumerate(stats):
        x = tile_x + (tile_w + gap) * i
        add_round(s, x, tile_y, tile_w, tile_h, SAGE_SOFT, corner=0.06)
        add_text(s, x + Inches(0.2), tile_y + Inches(0.18),
                 tile_w - Inches(0.4), Inches(0.85),
                 big, size=34, bold=True, color=FOREST, font=HEAD_FONT,
                 line_spacing=1.0)
        add_text(s, x + Inches(0.2), tile_y + Inches(0.95),
                 tile_w - Inches(0.4), Inches(0.32),
                 lbl, size=12, bold=True, color=INK, font=BODY_FONT,
                 line_spacing=1.0)
        add_text(s, x + Inches(0.2), tile_y + Inches(1.25),
                 tile_w - Inches(0.4), Inches(0.3),
                 sub, size=9, color=INK_SOFT, font=BODY_FONT, line_spacing=1.1)

    # Below: findings list
    add_kicker(s, Inches(0.5), Inches(3.7), Inches(12), "Findings")
    add_text(s, Inches(0.5), Inches(4.0), Inches(12), Inches(0.4),
             "Implementation matches the reports — with two clarifications",
             size=18, bold=True, color=INK, font=HEAD_FONT)
    bullets = [
        ("Target framework", "Tech Proposal §11 and Compliance Matrix §11.1 say .NET 8. The codebase targets net9.0 in every .csproj. Use .NET 9 on go-forward materials."),
        ("Vendor commitment VC-6", "SignalR notification hub is listed as an Enhancement (vendor commitment). It is already shipped — NotificationHub.cs + SignalRNotificationBroadcaster.cs are wired in WebAPI."),
        ("Domain entities, services, controllers, Razor pages", "Every path cited in Tech Proposal §3-§17 and in the Compliance Matrix rows is present at the stated location."),
        ("Security & RBAC", "AuthorizationPolicies.cs publishes internal.read/edit/manage/admin + external.investor; rate limiter enforces Free 60 / Plus 240 / Partner 600 req/min; ApiKey middleware bound to /api/v1/open-data/*; MFA mandatory for internal users."),
    ]
    by = Inches(4.6)
    for tag, body in bullets:
        add_rect(s, Inches(0.5), by + Inches(0.07), Emu(60000), Inches(0.45), GOLD)
        add_text(s, Inches(0.8), by, Inches(3.0), Inches(0.4),
                 tag, size=12, bold=True, color=FOREST, font=BODY_FONT,
                 line_spacing=1.2)
        add_text(s, Inches(3.9), by, Inches(8.9), Inches(0.5),
                 body, size=11, color=INK_SOFT, font=BODY_FONT, line_spacing=1.3)
        by += Inches(0.65)

    add_footer(s, prs, num, total)


def slide_architecture(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "Solution architecture", "Clean Architecture · CQRS · DDD")

    # Three horizontal bands: Hosts, Application+Domain, Infrastructure+Data
    band_x = Inches(0.5); band_w = Inches(12.3)
    # Hosts row
    row_y = Inches(1.85); row_h = Inches(1.25)
    hosts = [
        ("IRETP.Web", "Blazor Server portal\nPublic + Admin Razor pages"),
        ("IRETP.WebAPI", "Public + Investor API\nJWT + OIDC · rate-limited\nOpen Data /v1/open-data/*"),
        ("IRETP.AdminAPI", "Internal DLD operations\nEWRS · Escrow · Rating\nPrivate network only"),
    ]
    hw = (band_w - Inches(0.4)) / 3
    for i, (n, d) in enumerate(hosts):
        x = band_x + (hw + Inches(0.2)) * i
        add_round(s, x, row_y, hw, row_h, FOREST, corner=0.05)
        add_text(s, x + Inches(0.25), row_y + Inches(0.18), hw - Inches(0.5),
                 Inches(0.35), n, size=14, bold=True, color=PAPER,
                 font=HEAD_FONT, line_spacing=1.0)
        add_text(s, x + Inches(0.25), row_y + Inches(0.55), hw - Inches(0.5),
                 Inches(0.85), d, size=10.5, color=SAGE_SOFT,
                 font=BODY_FONT, line_spacing=1.3)

    # Middle band: Application / Domain
    mid_y = Inches(3.35); mid_h = Inches(1.65)
    add_round(s, band_x, mid_y, band_w, mid_h, SAGE_SOFT, corner=0.03)
    add_text(s, band_x + Inches(0.25), mid_y + Inches(0.15), band_w - Inches(0.5),
             Inches(0.35), "IRETP.Application  ·  CQRS via MediatR",
             size=13, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
    feats = [
        "AIAgent", "Alerts", "Analytics", "Audit", "Benchmark", "CMS",
        "Dashboard", "DeveloperRating", "EWRS", "Escrow", "Esg", "Export",
        "Greti", "Map", "Mortgage", "Ownership", "PriceIndex", "RentalIndex",
        "Transactions", "Auth",
    ]
    # Render features as small chip row
    chip_x = band_x + Inches(0.25); chip_y = mid_y + Inches(0.55)
    cw = Inches(1.16); ch = Inches(0.32); gx = Inches(0.06); gy = Inches(0.10)
    per_row = 10
    for i, f in enumerate(feats):
        r = i // per_row; c = i % per_row
        x = chip_x + (cw + gx) * c
        y = chip_y + (ch + gy) * r
        add_round(s, x, y, cw, ch, PAPER, corner=0.30, line=SAGE)
        add_text(s, x, y, cw, ch, f, size=10, bold=True, color=FOREST,
                 font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.0)
    add_text(s, band_x + Inches(0.25), mid_y + mid_h - Inches(0.30),
             band_w - Inches(0.5), Inches(0.25),
             "IRETP.Domain  ·  DDD entities + enums (Transaction, Project, EscrowAccount, RiskAlert, ScoringWeight, …)",
             size=10, italic=True, color=INK_SOFT, font=BODY_FONT, line_spacing=1.0)

    # Infrastructure band
    inf_y = Inches(5.20); inf_h = Inches(1.25)
    add_round(s, band_x, inf_y, band_w, inf_h, FOREST_DEEP, corner=0.03)
    add_text(s, band_x + Inches(0.25), inf_y + Inches(0.15), band_w - Inches(0.5),
             Inches(0.35), "IRETP.Infrastructure",
             size=13, bold=True, color=GOLD, font=HEAD_FONT, line_spacing=1.0)
    infra_items = [
        ("EF Core", "SQL Server / Postgres"),
        ("OneLake", "medallion bronze→silver→gold"),
        ("Hangfire", "9 recurring jobs"),
        ("AI Orchestrator", "RAG · multi-model fallback"),
        ("Identity", "JWT · OIDC · MFA · RBAC"),
        ("Notifications", "Email · SMS · SignalR"),
    ]
    iw = (band_w - Inches(0.5)) / 6
    for i, (n, d) in enumerate(infra_items):
        x = band_x + Inches(0.25) + iw * i + Inches(0.04) * i
        y = inf_y + Inches(0.55)
        add_text(s, x, y, iw - Inches(0.04), Inches(0.3),
                 n, size=11, bold=True, color=PAPER,
                 font=BODY_FONT, line_spacing=1.0)
        add_text(s, x, y + Inches(0.30), iw - Inches(0.04), Inches(0.34),
                 d, size=9, color=SAGE, font=BODY_FONT, line_spacing=1.2)

    add_text(s, band_x, Inches(6.65), band_w, Inches(0.3),
             "Strict inward-pointing dependencies — Domain has no Infrastructure references; Application has no EF Core.",
             size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER,
             font=BODY_FONT, line_spacing=1.0)

    add_footer(s, prs, num, total)


def slide_stack(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "Technology stack", "Runtime · data · jobs · observability")

    sections = [
        ("Runtime & app", FOREST, [
            (".NET 9.0",        "ASP.NET Core + Blazor Server"),
            ("MediatR",         "CQRS dispatcher in Application"),
            ("FluentValidation","handler input validation"),
            ("Serilog",         "structured logs · Seq / Azure Monitor"),
        ]),
        ("Data", SAGE, [
            ("SQL Server / Postgres", "OLTP, EF Core 9 no-tracking projections"),
            ("OneLake / Fabric",       "medallion lakehouse"),
            ("Redis",                  "distributed cache (commitment)"),
            ("Blob Storage",           "exports, PDF artefacts"),
        ]),
        ("Jobs & realtime", GOLD, [
            ("Hangfire",            "9 jobs · risk · KPI · escrow · benchmark"),
            ("SignalR",             "in-platform notification hub"),
            ("HTTP SMS / SMTP",     "alert delivery"),
            ("HmacUnsubscribe",     "RFC 8058 one-click unsubscribe"),
        ]),
        ("Security & ops", FOREST_DEEP, [
            ("JWT + OIDC",         "dual scheme; UAE Pass / Azure AD ready"),
            ("MFA (TOTP)",         "mandatory for internal users"),
            ("PartitionedRateLimiter","Free 60 / Plus 240 / Partner 600"),
            ("Health checks",      "/healthz/sla aggregates AI · KPI · notifications"),
        ]),
    ]
    grid_x = Inches(0.5); grid_y = Inches(1.85)
    cw = Inches(6.1); ch = Inches(2.45); gap = Inches(0.2)
    for i, (title, color, rows) in enumerate(sections):
        r = i // 2; c = i % 2
        x = grid_x + (cw + gap) * c
        y = grid_y + (ch + gap) * r
        add_round(s, x, y, cw, ch, PAPER, corner=0.04, line=DIVIDER)
        # Vertical accent
        add_rect(s, x, y, Inches(0.08), ch, color)
        add_text(s, x + Inches(0.25), y + Inches(0.15), cw - Inches(0.5),
                 Inches(0.35),
                 title, size=14, bold=True, color=color,
                 font=HEAD_FONT, line_spacing=1.0)
        for ri, (k, v) in enumerate(rows):
            row_y = y + Inches(0.6) + Inches(0.42) * ri
            add_text(s, x + Inches(0.25), row_y, Inches(1.95), Inches(0.32),
                     k, size=11, bold=True, color=INK,
                     font=BODY_FONT, line_spacing=1.1)
            add_text(s, x + Inches(2.3), row_y, cw - Inches(2.5),
                     Inches(0.32),
                     v, size=10.5, color=INK_SOFT, font=BODY_FONT,
                     line_spacing=1.2)

    add_footer(s, prs, num, total)


def slide_coverage_table(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "RFP coverage", "Compliance Matrix · summary roll-up")

    headers = ["RFP Section", "C", "CD", "E", "Total"]
    rows = [
        ("§3  External Public Portal",        "5",  "0", "0", "5"),
        ("§4  Interactive Analytics",         "3",  "0", "0", "3"),
        ("§5  Real Estate AI Agent",          "9",  "2", "1", "12"),
        ("§6  Investor Notifications",        "9",  "0", "1", "10"),
        ("§7  Multilingual",                  "7",  "2", "0", "9"),
        ("§8  EWRS & Escrow Monitoring",      "16", "0", "0", "16"),
        ("§9  Developer Performance",         "7",  "0", "0", "7"),
        ("§10 Non-Functional, Security & RBAC","9", "1", "1", "11"),
        ("§11 Technology & Hosting",          "2",  "4", "0", "6"),
        ("§12 Data Familiarisation",          "0",  "3", "0", "3"),
        ("§13 Phased Delivery Plan",          "1",  "3", "0", "4"),
        ("§19 Data Governance & Privacy",     "2",  "2", "0", "4"),
        ("§20 Additional Recommended",        "0",  "0", "8", "8"),
        ("Vendor Commitments (Value-Add)",    "0",  "0", "15", "15"),
    ]
    add_table(s, Inches(0.5), Inches(1.85), Inches(7.7), Inches(5.2),
              headers, rows, font_size=10.5,
              col_widths=[Inches(3.6), Inches(0.8), Inches(0.9), Inches(0.8), Inches(1.6)])

    # Legend column on the right
    lx = Inches(8.6); ly = Inches(1.85); lw = Inches(4.2)
    add_round(s, lx, ly, lw, Inches(5.2), SAGE_SOFT, corner=0.04)
    add_text(s, lx + Inches(0.3), ly + Inches(0.25), lw - Inches(0.6), Inches(0.4),
             "Status legend",
             size=14, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
    legend = [
        ("C",  FOREST,
         "Compliant today — the implementing component is live in the codebase. The Compliance Matrix points to the exact file or service."),
        ("CD", SAGE,
         "Compliant by delivery — capability is in the phase plan; foundations exist and the work is scoped, estimated and committed."),
        ("E",  GOLD,
         "Enhanced — goes beyond the literal requirement. Contractually offered at no additional charge in §20 and the Vendor Commitments group."),
    ]
    ey = ly + Inches(0.8)
    for tag, col, desc in legend:
        add_round(s, lx + Inches(0.3), ey, Inches(0.55), Inches(0.55), col, corner=0.5)
        add_text(s, lx + Inches(0.3), ey, Inches(0.55), Inches(0.55),
                 tag, size=18, bold=True, color=PAPER, font=HEAD_FONT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        add_text(s, lx + Inches(1.05), ey, lw - Inches(1.3), Inches(1.6),
                 desc, size=11, color=INK_SOFT, font=BODY_FONT, line_spacing=1.25)
        ey += Inches(1.4)

    add_footer(s, prs, num, total)


def slide_portal(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "RFP §3 · External Public Portal", "Public-facing transparency surface")

    rows = [
        ("§3.1  CMS",         "CmsContent + CmsContentVersion · versioned drafts; AdminAPI authoring + Public/CmsPreview.razor."),
        ("§3.2  Dashboard",   "Public/Dashboard.razor + DashboardController · KpiSnapshotCache refreshed every 15 min by KpiSnapshotRefreshService. FR-004 market charts (volume / type / top-5 zones / PSF dual-axis)."),
        ("§3.3  Transactions","Public/Transactions.razor · server-side filters · pagination · CSV/PDF export · URL persistence."),
        ("§3.4  GIS Map",     "Public/Map.razor + MapController · zone polygons + 511 project pins + 4 heatmap modes (volume / pricesqft / yield / ESG)."),
        ("§3.5  Price & Rental", "Public/PriceIndex.razor + Public/RentalIndex.razor · seeded 8 quarters per zone · yield calculator embedded."),
    ]
    by = Inches(1.85)
    for k, v in rows:
        # row container
        add_round(s, Inches(0.5), by, Inches(12.3), Inches(0.95), PAPER,
                  corner=0.04, line=DIVIDER)
        add_rect(s, Inches(0.5), by, Inches(0.08), Inches(0.95), FOREST)
        add_text(s, Inches(0.75), by + Inches(0.15), Inches(2.7), Inches(0.35),
                 k, size=13, bold=True, color=FOREST, font=HEAD_FONT,
                 line_spacing=1.0)
        add_text(s, Inches(0.75), by + Inches(0.5), Inches(11.9), Inches(0.4),
                 v, size=11, color=INK_SOFT, font=BODY_FONT, line_spacing=1.25)
        by += Inches(1.0)

    add_text(s, Inches(0.5), Inches(6.92), Inches(12.3), Inches(0.2),
             "Every page is bilingual (EN/AR with RTL), uses the shared <ModuleBanner>, <Skeleton>, <NetworkError>, <EmptyState> components.",
             size=9.5, italic=True, color=MUTED, font=BODY_FONT, line_spacing=1.0)
    add_footer(s, prs, num, total)


def slide_analytics(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "RFP §4 · Slice-and-Dice Analytics", "Interactive multi-dimension analysis")

    # Left: query flow
    add_round(s, Inches(0.5), Inches(1.85), Inches(7.0), Inches(5.0),
              PAPER, corner=0.04, line=DIVIDER)
    add_text(s, Inches(0.75), Inches(2.05), Inches(6.5), Inches(0.4),
             "Query path",
             size=14, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
    pipeline = [
        ("1", "Public/Analytics.razor", "User selects zone · type · developer · period · financing"),
        ("2", "AnalyticsController", "POST /api/analytics/execute"),
        ("3", "ExecuteAnalyticsQueryHandler", "MediatR · EF Core projection on gold layer"),
        ("4", "SavedAnalyticsView", "Persist filter set with 12-month signed share token"),
        ("5", "Features/Export", "CSV / PDF — locale-aware headers"),
    ]
    py = Inches(2.55)
    for n, title, sub in pipeline:
        add_round(s, Inches(0.75), py, Inches(0.4), Inches(0.4), GOLD, corner=0.5)
        add_text(s, Inches(0.75), py, Inches(0.4), Inches(0.4),
                 n, size=14, bold=True, color=FOREST_DEEP, font=HEAD_FONT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        add_text(s, Inches(1.3), py - Inches(0.02), Inches(6.0), Inches(0.3),
                 title, size=12, bold=True, color=INK, font=BODY_FONT, line_spacing=1.0)
        add_text(s, Inches(1.3), py + Inches(0.22), Inches(6.0), Inches(0.3),
                 sub, size=10.5, color=INK_SOFT, font=BODY_FONT, line_spacing=1.2)
        py += Inches(0.85)

    # Right: capability tiles
    rx = Inches(7.8); ry = Inches(1.85); rw = Inches(5.0); rh = Inches(5.0)
    add_round(s, rx, ry, rw, rh, SAGE_SOFT, corner=0.04)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), rw - Inches(0.6), Inches(0.4),
             "Capabilities",
             size=14, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
    add_bullets(s, rx + Inches(0.3), ry + Inches(0.7), rw - Inches(0.6),
                Inches(4.3), [
        "Multi-dimension filters across zone, property type, developer, period and financing method.",
        "Saved Views capped at 12 per user; shareable URL signed with HMAC and 12-month expiry.",
        "Embed mode (AnalyticsEmbed.razor + EmbedLayout) for partner sites and DLD micro-sites.",
        "Charts respect locale — axes, tooltips, AR digits, RTL flip handled by chart-interop.js defaults.",
        "P95 chart render ≤ 2 s achieved through KpiSnapshotCache and EF AsNoTracking projections.",
    ], size=12, color=INK_SOFT)

    add_footer(s, prs, num, total)


def slide_ai(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "RFP §5 · Real Estate AI Agent", "Model-agnostic orchestration · RAG · guardrails")

    # Five horizontal pillars
    pillars = [
        ("Orchestrator",
         "AIOrchestrator routes per-task between provider tiers (nav / analytics / primary / secondary). Warm fallback on provider failure; per-model latency surfaced in Admin/AIModels.razor."),
        ("RAG",
         "BuildRagContextAsync indexes transactions, project registry, developer profiles, RERA regulations and DLD service pages. Every reply grounded in DLD data — mandatory and model-agnostic."),
        ("Guardrails",
         "KeywordAdvisoryGuardrail + regex layer blocks forward-looking price predictions, regulated investment advice. Bilingual refusal copy; original answer logged for audit."),
        ("Memory",
         "UserAiMemory + UserConsentService — top-5 zones and top-3 topics prepended after explicit opt-in from Account.razor. Erasure honoured on consent revoke."),
        ("Multilingual",
         "EN + AR conversational support in P1 via locale-aware prompts. ZH · RU · UR · FR · HI · DE extended in P4 via same orchestrator with per-language model selection."),
    ]
    px = Inches(0.5); py = Inches(1.85); pw = Inches(2.42); ph = Inches(4.4); gap = Inches(0.12)
    for i, (name, body) in enumerate(pillars):
        x = px + (pw + gap) * i
        add_round(s, x, py, pw, ph, PAPER, corner=0.04, line=DIVIDER)
        add_rect(s, x, py, pw, Inches(0.08), GOLD)
        add_text(s, x + Inches(0.18), py + Inches(0.3), pw - Inches(0.36),
                 Inches(0.55), name,
                 size=15, bold=True, color=FOREST, font=HEAD_FONT,
                 line_spacing=1.05)
        add_text(s, x + Inches(0.18), py + Inches(0.9), pw - Inches(0.36),
                 ph - Inches(1.1),
                 body, size=10.5, color=INK_SOFT, font=BODY_FONT, line_spacing=1.4)

    # Accuracy + UAE residency footer block
    fy = Inches(6.4); fw = Inches(12.3)
    add_round(s, Inches(0.5), fy, fw, Inches(0.7), FOREST_DEEP, corner=0.04)
    add_text(s, Inches(0.7), fy + Inches(0.1), fw - Inches(0.4), Inches(0.25),
             "ACCURACY · RESIDENCY",
             size=10, bold=True, color=GOLD, font=BODY_FONT, line_spacing=1.0)
    add_text(s, Inches(0.7), fy + Inches(0.32), fw - Inches(0.4), Inches(0.4),
             "AiAccuracyHarness runs a 14-question seed catalog (EN+AR, data + service-nav + adversarial). All inference routed through UAE-resident endpoints.",
             size=11.5, color=PAPER, font=BODY_FONT, line_spacing=1.2)

    add_footer(s, prs, num, total)


def slide_alerts(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "RFP §6 · Investor Notifications", "Six alert types · three channels · contracted SLAs")

    # Left: alert types
    add_round(s, Inches(0.5), Inches(1.85), Inches(6.1), Inches(5.0),
              PAPER, corner=0.04, line=DIVIDER)
    add_text(s, Inches(0.75), Inches(2.05), Inches(5.6), Inches(0.4),
             "Alert types (§6.1)",
             size=14, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
    types = [
        ("Price Movement", "Per-user 1–25 % threshold"),
        ("New Project Launch", "Watchlist scoped"),
        ("Watchlist Project Status", "Construction / certification change"),
        ("Rental Yield Threshold", "Below / above user band"),
        ("Periodic Market Digest", "Weekly · monthly · HTML w/ embedded chart"),
        ("Regulation / Policy Update", "From CmsContent + RegulatoryViolation"),
    ]
    ty = Inches(2.55)
    for n, d in types:
        add_rect(s, Inches(0.75), ty + Inches(0.16), Emu(60000), Inches(0.25), GOLD)
        add_text(s, Inches(1.05), ty, Inches(2.7), Inches(0.35),
                 n, size=11.5, bold=True, color=INK, font=BODY_FONT, line_spacing=1.1)
        add_text(s, Inches(3.6), ty, Inches(3.0), Inches(0.35),
                 d, size=10.5, color=INK_SOFT, font=BODY_FONT, line_spacing=1.2)
        ty += Inches(0.55)

    # Right: channels + SLA + tech
    rx = Inches(6.8); ry = Inches(1.85); rw = Inches(6.0)
    add_round(s, rx, ry, rw, Inches(2.0), SAGE_SOFT, corner=0.04)
    add_text(s, rx + Inches(0.3), ry + Inches(0.18), rw - Inches(0.6),
             Inches(0.4), "Channels (§6.2)",
             size=14, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
    ch = [
        ("Email", "SmtpEmailSender · branded HTML + chart · HMAC unsubscribe (RFC 8058)", "≤ 5 min"),
        ("SMS",   "HttpSmsSender · ≤ 160 chars · shortened link",                     "≤ 3 min"),
        ("In-platform", "Notification entity · SignalR hub · bell widget unread badge", "instant"),
    ]
    chy = ry + Inches(0.6)
    for n, d, sla in ch:
        add_text(s, rx + Inches(0.3), chy, Inches(1.5), Inches(0.3),
                 n, size=12, bold=True, color=INK, font=BODY_FONT, line_spacing=1.0)
        add_text(s, rx + Inches(1.6), chy, Inches(3.4), Inches(0.3),
                 d, size=10, color=INK_SOFT, font=BODY_FONT, line_spacing=1.15)
        add_text(s, rx + rw - Inches(0.95), chy, Inches(0.7), Inches(0.3),
                 sla, size=11, bold=True, color=FOREST,
                 font=BODY_FONT, align=PP_ALIGN.RIGHT, line_spacing=1.0)
        chy += Inches(0.45)

    add_round(s, rx, ry + Inches(2.2), rw, Inches(2.8), PAPER, corner=0.04, line=DIVIDER)
    add_text(s, rx + Inches(0.3), ry + Inches(2.4), rw - Inches(0.6), Inches(0.4),
             "Pipeline",
             size=14, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
    add_bullets(s, rx + Inches(0.3), ry + Inches(2.9), rw - Inches(0.6),
                Inches(2.2), [
        "InvestorAlertEvaluator runs as a Hangfire job, joining InvestorAlert + WatchlistItem + signal sources.",
        "AlertDeliveryService dispatches via NotificationRecipientResolver, honouring channel preferences and quiet hours.",
        "NotificationTemplates centralises HTML / SMS copy; List-Unsubscribe + List-Unsubscribe-Post headers on every send.",
        "NotificationSlaHealthCheck reports P95 vs RFP budget on /healthz/sla.",
    ], size=11, color=INK_SOFT)

    add_footer(s, prs, num, total)


def slide_ewrs(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "RFP §8 · Early Warning Risk System",
               "Ten configurable indicators · L1–L4 escalation")

    indicators = [
        ("Project Delivery Delay",         "warning / critical"),
        ("Escrow Shortfall",               "warning / critical"),
        ("Construction Activity Suspension","30 / 60 day threshold"),
        ("Transaction Volume Decline",     "vs. 12-mo rolling avg"),
        ("Zone Price Decline",             "PSF QoQ, min 5 tx"),
        ("Developer Score Deterioration",  "QoQ; <40 escalates High"),
        ("High-Risk Project Concentration","≥3 active w/ open High"),
        ("Severe Regulatory Violation",    "Critical RERA in 90d"),
        ("AI Anomaly Flag",                "TimeSeriesAnalyzer z-score"),
        ("Data-Feed Anomaly",              "publication held until ack"),
    ]
    # Two-column indicator list
    px = Inches(0.5); py = Inches(1.85); pw = Inches(7.6)
    add_round(s, px, py, pw, Inches(5.1), PAPER, corner=0.04, line=DIVIDER)
    add_text(s, px + Inches(0.3), py + Inches(0.2), pw - Inches(0.6),
             Inches(0.4),
             "Ten risk indicators — all thresholds configurable from EwrsDashboard",
             size=12.5, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
    for i, (n, d) in enumerate(indicators):
        col = i % 2; row = i // 2
        x = px + Inches(0.3) + col * Inches(3.5)
        y = py + Inches(0.75) + row * Inches(0.85)
        # Bullet square
        add_rect(s, x, y + Inches(0.08), Inches(0.15), Inches(0.15), GOLD)
        add_text(s, x + Inches(0.3), y - Inches(0.05), Inches(3.2),
                 Inches(0.35), n,
                 size=11.5, bold=True, color=INK, font=BODY_FONT, line_spacing=1.0)
        add_text(s, x + Inches(0.3), y + Inches(0.22), Inches(3.2),
                 Inches(0.35), d,
                 size=10, color=INK_SOFT, font=BODY_FONT, line_spacing=1.1)

    # Escalation column on the right
    ex = Inches(8.3); ey = Inches(1.85); ew = Inches(4.5)
    add_round(s, ex, ey, ew, Inches(5.1), SAGE_SOFT, corner=0.04)
    add_text(s, ex + Inches(0.3), ey + Inches(0.2), ew - Inches(0.6), Inches(0.4),
             "Escalation (§8.2)",
             size=13, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
    levels = [
        ("L1", "Project Officer", "Email + Platform"),
        ("L2", "Section Manager + Director", "Email + SMS + Platform"),
        ("L3", "DG + Deputies", "All channels + SMS"),
        ("L4", "DG + Regulator", "All channels + Regulator notice"),
    ]
    ly = ey + Inches(0.7)
    for lvl, who, chn in levels:
        add_round(s, ex + Inches(0.3), ly, Inches(0.65), Inches(0.7), FOREST, corner=0.15)
        add_text(s, ex + Inches(0.3), ly, Inches(0.65), Inches(0.7),
                 lvl, size=18, bold=True, color=GOLD, font=HEAD_FONT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        add_text(s, ex + Inches(1.1), ly + Inches(0.04), ew - Inches(1.4),
                 Inches(0.35),
                 who, size=11.5, bold=True, color=INK, font=BODY_FONT,
                 line_spacing=1.0)
        add_text(s, ex + Inches(1.1), ly + Inches(0.35), ew - Inches(1.4),
                 Inches(0.35),
                 chn, size=10, color=INK_SOFT, font=BODY_FONT, line_spacing=1.1)
        ly += Inches(0.92)

    add_text(s, ex + Inches(0.3), ly + Inches(0.1), ew - Inches(0.6), Inches(0.5),
             "AlertEscalationService auto-promotes alerts whose SLA window has elapsed (Hangfire every 5 min). Every threshold edit is audited.",
             size=10, italic=True, color=INK_SOFT, font=BODY_FONT, line_spacing=1.25)

    add_footer(s, prs, num, total)


def slide_developer(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "RFP §9 + §8.4 · Developer Rating & Escrow",
               "Transparent composite score · real-time escrow oversight")

    # Left: 8 scoring criteria
    add_round(s, Inches(0.5), Inches(1.85), Inches(6.1), Inches(5.0),
              PAPER, corner=0.04, line=DIVIDER)
    add_text(s, Inches(0.75), Inches(2.05), Inches(5.6), Inches(0.4),
             "Eight scoring criteria (§9.1.1)",
             size=14, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
    criteria = [
        "On-Time Project Delivery Rate",
        "Unit Sales Completion Rate",
        "Escrow Account Health Score",
        "Regulatory Compliance Record",
        "Customer Complaints",
        "Construction Quality",
        "Financial Stability",
        "Innovation & Sustainability",
    ]
    cy = Inches(2.55)
    for i, c in enumerate(criteria):
        add_rect(s, Inches(0.85), cy + Inches(0.13), Emu(50000), Inches(0.25), GOLD)
        add_text(s, Inches(1.1), cy, Inches(3.6), Inches(0.35),
                 c, size=11.5, color=INK, font=BODY_FONT, line_spacing=1.1)
        add_text(s, Inches(5.0), cy, Inches(1.5), Inches(0.35),
                 "ScoringWeight", size=10, italic=True, color=MUTED,
                 font=BODY_FONT, align=PP_ALIGN.RIGHT, line_spacing=1.1)
        cy += Inches(0.43)
    add_text(s, Inches(0.75), Inches(6.25), Inches(5.6), Inches(0.5),
             "All weights adjustable from Admin/DeveloperRating.razor · every change audited with timestamp + approver identity.",
             size=10, italic=True, color=INK_SOFT, font=BODY_FONT, line_spacing=1.25)

    # Right: Escrow trio
    rx = Inches(6.8); ry = Inches(1.85); rw = Inches(6.0)
    add_round(s, rx, ry, rw, Inches(5.0), SAGE_SOFT, corner=0.04)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), rw - Inches(0.6), Inches(0.4),
             "Escrow monitoring (§8.4)",
             size=14, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
    esc = [
        ("ESC-001  Real-time dashboard",
         "Admin/EscrowMonitoring + EscrowDetail · adequacy ratio + Adequate / Warning / Critical badges."),
        ("ESC-002  Immutable audit log",
         "EscrowTransaction is append-only · PDF export via InvestorScorecardPdfService · L3 alert on unauthorised discrepancy."),
        ("ESC-003  Monthly health report",
         "EscrowHealthReportService Hangfire job · per-project PDF generated month-end · emailed to assigned officer within 24 h."),
    ]
    ey = ry + Inches(0.7)
    for n, d in esc:
        add_text(s, rx + Inches(0.3), ey, rw - Inches(0.6), Inches(0.35),
                 n, size=12, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
        add_text(s, rx + Inches(0.3), ey + Inches(0.35), rw - Inches(0.6),
                 Inches(0.95), d,
                 size=11, color=INK_SOFT, font=BODY_FONT, line_spacing=1.3)
        ey += Inches(1.3)
    add_footer(s, prs, num, total)


def slide_multilingual(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "RFP §7 · Multilingual · Currency",
               "EN + AR first-class · 6-language extension in P4")

    # Top: language matrix
    headers = ["Channel", "Phase 1 (EN + AR)", "Phase 4 extension"]
    rows = [
        ("Internal Platform UI",        "C — MainLayout · per-page CSS · resources",        "—"),
        ("External Portal UI",          "C — Public/* pages · RTL · digits",                "ZH · RU · UR · FR · HI · DE"),
        ("AI Agent (chat)",             "C — locale-aware prompts",                         "Per-language model selection"),
        ("Investor alerts (email/SMS)", "C — NotificationTemplates locale fields",          "Extended templates per language"),
        ("CMS content",                 "C — CmsContent locale fields",                     "Authors add languages on demand"),
        ("Chart / export labels",       "C — chart-interop.js + Features/Export",           "Inherits page locale"),
    ]
    add_table(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(2.9),
              headers, rows, font_size=10.5,
              col_widths=[Inches(3.6), Inches(5.0), Inches(3.7)])

    # Bottom: two side notes
    by = Inches(4.95); bw = Inches(6.0); bh = Inches(1.95)
    add_round(s, Inches(0.5), by, bw, bh, PAPER, corner=0.04, line=DIVIDER)
    add_rect(s, Inches(0.5), by, Inches(0.08), bh, GOLD)
    add_text(s, Inches(0.75), by + Inches(0.15), bw - Inches(0.4), Inches(0.4),
             "RTL & i18n",
             size=14, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
    add_text(s, Inches(0.75), by + Inches(0.55), bw - Inches(0.4),
             bh - Inches(0.7),
             "Layout mirroring for AR / UR is baked into MainLayout and the design-system CSS. "
             "All translatable copy flows through resource files; no hardcoded strings in Razor markup. "
             "AR validation of zone / project / developer names is signed off against DLD official records before go-live.",
             size=11, color=INK_SOFT, font=BODY_FONT, line_spacing=1.35)

    bx2 = Inches(6.8)
    add_round(s, bx2, by, bw, bh, PAPER, corner=0.04, line=DIVIDER)
    add_rect(s, bx2, by, Inches(0.08), bh, GOLD)
    add_text(s, bx2 + Inches(0.25), by + Inches(0.15), bw - Inches(0.4),
             Inches(0.4),
             "Currency switcher",
             size=14, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
    add_text(s, bx2 + Inches(0.25), by + Inches(0.55), bw - Inches(0.4),
             bh - Inches(0.7),
             "CurrencyController + CurrencyRate entity + CurrencyRatesRefreshService feed a header switcher that converts AED to the selected currency. "
             "User language and currency preferences are persisted on ApplicationUser and survive navigation and re-login.",
             size=11, color=INK_SOFT, font=BODY_FONT, line_spacing=1.35)

    add_footer(s, prs, num, total)


def slide_security(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "RFP §10 + DESC ISR v3 · Security & RBAC",
               "Policy catalog · 35 ISR controls mapped")

    # Left: RBAC policies
    add_round(s, Inches(0.5), Inches(1.85), Inches(6.1), Inches(5.0),
              PAPER, corner=0.04, line=DIVIDER)
    add_text(s, Inches(0.75), Inches(2.05), Inches(5.6), Inches(0.4),
             "RBAC policy catalog",
             size=14, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
    pols = [
        ("internal.read",    "Viewer + above"),
        ("internal.edit",    "Operator + above"),
        ("internal.manage",  "Supervisor + above"),
        ("internal.admin",   "SystemAdministrator only"),
        ("external.investor","RegisteredInvestor + above"),
    ]
    py = Inches(2.6)
    for k, v in pols:
        add_round(s, Inches(0.85), py, Inches(2.6), Inches(0.4), SAGE_SOFT, corner=0.3)
        add_text(s, Inches(0.85), py, Inches(2.6), Inches(0.4),
                 k, size=11.5, bold=True, color=FOREST, font=BODY_FONT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        add_text(s, Inches(3.7), py + Inches(0.07), Inches(2.7), Inches(0.3),
                 v, size=11, color=INK_SOFT, font=BODY_FONT, line_spacing=1.0)
        py += Inches(0.55)

    add_text(s, Inches(0.75), py + Inches(0.05), Inches(5.6), Inches(0.8),
             "Cumulative roles — a Supervisor satisfies internal.read, internal.edit and internal.manage. "
             "Every controller is decorated with [Authorize(Policy = …)]. "
             "Rate limiter enforces Free 60 · Plus 240 · Partner 600 req/min keyed on X-Api-Key.",
             size=10, italic=True, color=INK_SOFT, font=BODY_FONT, line_spacing=1.3)

    # Right: DESC ISR posture
    rx = Inches(6.8); ry = Inches(1.85); rw = Inches(6.0)
    add_round(s, rx, ry, rw, Inches(5.0), SAGE_SOFT, corner=0.04)
    add_text(s, rx + Inches(0.3), ry + Inches(0.2), rw - Inches(0.6), Inches(0.4),
             "DESC ISR v3 control posture",
             size=14, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
    # Three big stat tiles in a row
    big = [("30", "Implemented", FOREST), ("3", "Partial", GOLD), ("4", "Planned", SAGE)]
    bw = (rw - Inches(0.8)) / 3
    for i, (n, lbl, col) in enumerate(big):
        x = rx + Inches(0.3) + (bw + Inches(0.1)) * i
        add_round(s, x, ry + Inches(0.7), bw, Inches(1.3), PAPER, corner=0.04, line=DIVIDER)
        add_text(s, x, ry + Inches(0.78), bw, Inches(0.7),
                 n, size=34, bold=True, color=col, font=HEAD_FONT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        add_text(s, x, ry + Inches(1.55), bw, Inches(0.3),
                 lbl, size=11, bold=True, color=INK, font=BODY_FONT,
                 align=PP_ALIGN.CENTER, line_spacing=1.0)
    add_text(s, rx + Inches(0.3), ry + Inches(2.15), rw - Inches(0.6),
             Inches(0.4),
             "Ten domains · 37 controls in Appendix A",
             size=11, italic=True, color=INK_SOFT, font=BODY_FONT,
             align=PP_ALIGN.CENTER, line_spacing=1.0)

    add_bullets(s, rx + Inches(0.3), ry + Inches(2.7), rw - Inches(0.6),
                Inches(2.2), [
        "JWT (HMAC-SHA256, zero clock-skew) + dual-scheme OIDC (UAE Pass / Azure AD).",
        "Append-only AuditLog enforced at DbContext; immutability test in xUnit suite.",
        "TLS 1.2+ everywhere · WAF · private endpoints · CORS allow-list (not wildcard).",
        "DESC-authorised VAPT scheduled before Phase 3 acceptance.",
    ], size=10.5, color=INK_SOFT)

    add_footer(s, prs, num, total)


def slide_hosting(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "RFP §11 · Hosting & Disaster Recovery",
               "Three options · UAE-resident · ≤ 4 h RPO")

    options = [
        ("Option A", "Azure UAE North",
         "Recommended for fastest TTL with built-in BCDR.",
         "AED 335.5 K", "AED 629.2 K"),
        ("Option B", "On-Premises (DLD-approved DC)",
         "Preferred for full data-sovereignty assurance.",
         "AED 1.41 M", "AED 686.2 K"),
        ("Option C", "Hybrid",
         "Compute on-prem · lakehouse in Azure UAE · site-to-site link.",
         "AED 793.5 K", "AED 654.2 K"),
    ]
    cx = Inches(0.5); cy = Inches(1.85); cw = Inches(4.05); ch = Inches(3.4); gap = Inches(0.15)
    for i, (tag, name, line, capex, opex) in enumerate(options):
        x = cx + (cw + gap) * i
        add_round(s, x, cy, cw, ch, PAPER, corner=0.04, line=DIVIDER)
        add_rect(s, x, cy, cw, Inches(0.08), GOLD)
        add_text(s, x + Inches(0.25), cy + Inches(0.25), cw - Inches(0.5),
                 Inches(0.35),
                 tag, size=10, bold=True, color=GOLD, font=BODY_FONT, line_spacing=1.0)
        add_text(s, x + Inches(0.25), cy + Inches(0.55), cw - Inches(0.5),
                 Inches(0.65),
                 name, size=17, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.1)
        add_text(s, x + Inches(0.25), cy + Inches(1.30), cw - Inches(0.5),
                 Inches(0.85),
                 line, size=11, color=INK_SOFT, font=BODY_FONT, line_spacing=1.3)
        # CAPEX / OPEX strip
        sy = cy + Inches(2.30)
        add_rect(s, x + Inches(0.25), sy, cw - Inches(0.5), Inches(0.04), DIVIDER)
        add_text(s, x + Inches(0.25), sy + Inches(0.10), Inches(1.6), Inches(0.3),
                 "CAPEX", size=9, bold=True, color=MUTED, font=BODY_FONT, line_spacing=1.0)
        add_text(s, x + Inches(0.25), sy + Inches(0.35), Inches(1.8), Inches(0.35),
                 capex, size=13, bold=True, color=INK, font=BODY_FONT, line_spacing=1.0)
        add_text(s, x + cw - Inches(1.85), sy + Inches(0.10), Inches(1.6), Inches(0.3),
                 "ANNUAL OPEX", size=9, bold=True, color=MUTED,
                 font=BODY_FONT, line_spacing=1.0)
        add_text(s, x + cw - Inches(1.85), sy + Inches(0.35), Inches(1.6),
                 Inches(0.35),
                 opex, size=13, bold=True, color=INK, font=BODY_FONT, line_spacing=1.0)

    # DR strip
    dy = Inches(5.55); dw = Inches(12.3)
    add_round(s, Inches(0.5), dy, dw, Inches(1.4), SAGE_SOFT, corner=0.04)
    add_text(s, Inches(0.7), dy + Inches(0.15), dw - Inches(0.4), Inches(0.35),
             "Mandatory hosting controls (§11.3) — apply to every option",
             size=12, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
    add_text(s, Inches(0.7), dy + Inches(0.55), dw - Inches(0.4), Inches(0.85),
             "TLS 1.2+ end-to-end  ·  segregated networks for WebAPI vs AdminAPI  ·  WAF in front of WebAPI  ·  private endpoints for SQL · Redis · Storage\n"
             "BCDR with cross-AZ replicas  ·  RTO ≤ 8 h  ·  RPO ≤ 4 h  ·  geo-redundant backups  ·  slot-swap zero-downtime release",
             size=11, color=INK_SOFT, font=BODY_FONT, line_spacing=1.5)

    add_footer(s, prs, num, total)


def slide_data(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "Data architecture", "Medallion lakehouse  ·  15-minute freshness")

    layers = [
        ("Bronze", FOREST_DEEP,
         "Raw ingestion",
         "DLD source-of-truth registries, escrow bank feeds (signed payload + checksum), RERA regulatory feeds, ESG certifications. No transformations."),
        ("Silver", FOREST,
         "Conformed & cleansed",
         "Data-quality rules (Great Expectations / dbt tests). Anomalies pause publication and raise a steward acknowledgement task."),
        ("Gold",   SAGE,
         "Business-ready",
         "EF Core 9 reads the gold layer for transactional UI. Hot KPIs synchronised into the OLTP store every 15 min by Hangfire."),
    ]
    by = Inches(1.85)
    for n, col, sub, desc in layers:
        h = Inches(1.55)
        add_round(s, Inches(0.5), by, Inches(12.3), h, PAPER, corner=0.04, line=DIVIDER)
        add_rect(s, Inches(0.5), by, Inches(2.0), h, col)
        add_text(s, Inches(0.5), by, Inches(2.0), h,
                 n, size=28, bold=True, color=GOLD, font=HEAD_FONT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        add_text(s, Inches(2.7), by + Inches(0.2), Inches(9.8), Inches(0.4),
                 sub, size=14, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
        add_text(s, Inches(2.7), by + Inches(0.6), Inches(9.8), h - Inches(0.7),
                 desc, size=11.5, color=INK_SOFT, font=BODY_FONT, line_spacing=1.4)
        by += h + Inches(0.15)

    add_text(s, Inches(0.5), by + Inches(0.1), Inches(12.3), Inches(0.4),
             "KpiSnapshotRefreshService · BenchmarkRefreshService · CurrencyRatesRefreshService — three Hangfire workers keep the gold-layer warm.",
             size=10, italic=True, color=MUTED, font=BODY_FONT,
             align=PP_ALIGN.CENTER, line_spacing=1.0)

    add_footer(s, prs, num, total)


def slide_phases(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "Phased delivery", "Four phases · 4 months · acceptance-gated")

    # All phase pills use FOREST so GOLD title text reads consistently.
    # The accent color travels as a thin left rule on the row instead.
    phases = [
        ("Phase 1", "Month 1", FOREST,
         "External Public Portal (CMS · Dashboard · Transactions · Map · Indices · Analytics · AI Agent EN/AR) · Open Data API + Developer Portal · Investor Notifications"),
        ("Phase 2", "Month 2", SAGE,
         "Investor Watchlist v2 · ESG indicators · GRETI sub-index closure · performance hardening (Redis · OpenTelemetry · load-test sign-off)"),
        ("Phase 3", "Month 3", GOLD,
         "Internal Platform — EWRS · Developer Rating · Escrow Monitoring · Beneficial Ownership · Mortgage analytics · DESC-authorised VAPT"),
        ("Phase 4", "Month 4", FOREST_DEEP,
         "Multilingual extension (ZH · RU · UR · FR · HI · DE) · AI Agent extended languages · advanced analytics · final UAT and acceptance"),
    ]
    by = Inches(1.8)
    for n, m, accent, body in phases:
        h = Inches(1.15); pw = Inches(12.3)
        add_round(s, Inches(0.5), by, pw, h, PAPER, corner=0.04, line=DIVIDER)
        # Accent rule on the left (varies per phase — this is where the
        # phase identity colour lives, not behind the title text)
        add_rect(s, Inches(0.5), by, Inches(0.1), h, accent)
        # Title pill — uniform FOREST so GOLD text always reads
        add_round(s, Inches(0.75), by + Inches(0.2), Inches(1.5), Inches(0.78),
                  FOREST, corner=0.15)
        add_text(s, Inches(0.75), by + Inches(0.2), Inches(1.5), Inches(0.38),
                 n, size=13, bold=True, color=GOLD, font=HEAD_FONT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        add_text(s, Inches(0.75), by + Inches(0.55), Inches(1.5), Inches(0.4),
                 m, size=10, color=PAPER, font=BODY_FONT,
                 align=PP_ALIGN.CENTER, line_spacing=1.0)
        add_text(s, Inches(2.5), by + Inches(0.23), pw - Inches(2.2),
                 Inches(0.85),
                 body, size=11.5, color=INK_SOFT, font=BODY_FONT,
                 line_spacing=1.4, anchor=MSO_ANCHOR.MIDDLE)
        by += h + Inches(0.1)

    add_text(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.2),
             "Each month exits with a §17.3 UAT bar — functional · performance · security · accessibility · language/RTL. Possible only because the codebase already covers ~88% of the RFP scope.",
             size=9.5, italic=True, color=MUTED, font=BODY_FONT,
             align=PP_ALIGN.CENTER, line_spacing=1.0)

    add_footer(s, prs, num, total)


def slide_responsibility(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "Responsibility matrix",
               "DLD · Vendor · Partner — role per programme activity")

    headers = ["Activity", "DLD", "Vendor", "Partner"]
    # Cell values use single-letter codes for clean reading
    rows = [
        ("Programme governance & steering",                  "A", "L", "S"),
        ("Solution architecture & design",                   "I", "L", "S"),
        ("Application development (.NET 9 · Blazor · CQRS)", "I", "L", "—"),
        ("Infrastructure provisioning & hosting",            "A", "S", "L"),
        ("Data discovery, migration & lakehouse build",      "S", "L", "S"),
        ("AI / RAG pipeline & guardrails",                   "A", "L", "I"),
        ("Security hardening · DESC ISR v3 · VAPT",          "A", "L", "S"),
        ("QA · test execution · UAT facilitation",           "S", "L", "S"),
        ("Acceptance sign-off",                              "A", "I", "I"),
        ("Post-go-live operations, support & SLA",           "I", "L", "S"),
    ]
    add_table(s, Inches(0.5), Inches(1.85), Inches(8.6), Inches(4.8),
              headers, rows, font_size=11,
              col_widths=[Inches(5.6), Inches(1.0), Inches(1.0), Inches(1.0)])

    # Centre-align the L/S/A/I cells for readability
    tbl_shape = s.shapes[-1]
    tbl = tbl_shape.table
    for ri in range(1, len(rows) + 1):
        for ci in (1, 2, 3):
            cell = tbl.cell(ri, ci)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER

    # Legend column on the right
    lx = Inches(9.3); ly = Inches(1.85); lw = Inches(3.5)
    add_round(s, lx, ly, lw, Inches(4.8), SAGE_SOFT, corner=0.04)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), lw - Inches(0.6), Inches(0.4),
             "Legend",
             size=14, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
    legend = [
        ("L", FOREST,       "Lead — primary executor and owner of the deliverable"),
        ("S", SAGE,         "Support — contributes effort, expertise or sign-off review"),
        ("A", GOLD,         "Approver — formal sign-off authority"),
        ("I", FOREST_DEEP,  "Informed — kept in the loop, no decision authority"),
    ]
    ey = ly + Inches(0.7)
    for tag, col, desc in legend:
        add_round(s, lx + Inches(0.3), ey, Inches(0.5), Inches(0.5), col, corner=0.5)
        add_text(s, lx + Inches(0.3), ey, Inches(0.5), Inches(0.5),
                 tag, size=15, bold=True, color=PAPER, font=HEAD_FONT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        add_text(s, lx + Inches(0.95), ey + Inches(0.02), lw - Inches(1.2),
                 Inches(1.0),
                 desc, size=10, color=INK_SOFT, font=BODY_FONT, line_spacing=1.25)
        ey += Inches(1.0)

    add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.3),
             "VAPT is executed by a DESC-approved third-party firm engaged by the Vendor; results are reviewed jointly with DLD before Phase 3 sign-off.",
             size=10, italic=True, color=MUTED, font=BODY_FONT, line_spacing=1.0)

    add_footer(s, prs, num, total)


def slide_powerbi_compare(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "Visualisation layer choice",
               "Power BI Embedded vs. proposed Blazor + Chart.js")

    headers = ["Dimension", "Power BI Embedded (A-series)", "Proposed: Blazor + Chart.js"]
    rows = [
        ("Year-1 licence cost (public users)",
         "AED 260 K – 1 M+ (A4–A6 capacity)",
         "AED 0 — MIT / open-source"),
        ("Anonymous public access",
         "App-owns-data token broker + service principal",
         "Native — anonymous by default"),
        ("UI / UX customisation depth",
         "Canvas-bound; theme JSON only; custom visuals need TypeScript",
         "Full HTML/CSS/JS; same design tokens as rest of IRETP"),
        ("RTL (Arabic) parity",
         "Incomplete on several visual types",
         "Layout-level; every visual inherits page direction"),
        ("Concurrent users per AED 100 K",
         "~50–100 (A5 capacity share)",
         "~500–1,500 (App Service Plan share)"),
        ("P95 chart render",
         "5–15 s cold without pre-warming",
         "≤ 2 s via KpiSnapshotCache + no-tracking projections"),
        ("Source control · test automation",
         "Binary .pbix; weak git diffs; minimal test tooling",
         "Razor + JS in .NET solution; xUnit + Playwright + CI"),
        ("DLD IP ownership (RFP §19.1)",
         "Microsoft-owned service surface",
         "Code vests in DLD on payment"),
    ]
    add_table(s, Inches(0.5), Inches(1.75), Inches(12.3), Inches(4.4),
              headers, rows, font_size=10.5,
              col_widths=[Inches(3.8), Inches(4.2), Inches(4.3)])

    # Verdict strip — compact, single line above the footer
    vy = Inches(6.30); vw = Inches(12.3)
    add_round(s, Inches(0.5), vy, vw, Inches(0.75), FOREST_DEEP, corner=0.04)
    add_text(s, Inches(0.7), vy + Inches(0.10), Inches(2.3), Inches(0.25),
             "RECOMMENDATION",
             size=10, bold=True, color=GOLD, font=BODY_FONT, line_spacing=1.0)
    add_text(s, Inches(0.7), vy + Inches(0.35), vw - Inches(0.4), Inches(0.35),
             "Retain Blazor + Chart.js for the External Portal. Power BI stays optional for internal DLD analyst self-service (Pro / PPU tenancy — no A-series capacity).",
             size=11, color=PAPER, font=BODY_FONT, line_spacing=1.2)

    add_footer(s, prs, num, total)


def slide_quality(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "Quality posture", "Operational schedule · forward QA plan · observability")

    # Top: 4 stat tiles — no test claims
    tiles = [
        ("9",            "Hangfire recurring jobs"),
        ("P95 ≤ 2 s",    "chart render target"),
        ("P95 ≤ 15 s",   "AI agent response target"),
        ("/healthz/sla", "aggregated SLO health"),
    ]
    tx = Inches(0.5); ty = Inches(1.85); tw = Inches(3.07); th = Inches(1.4); gap = Inches(0.13)
    for i, (big, lbl) in enumerate(tiles):
        x = tx + (tw + gap) * i
        add_round(s, x, ty, tw, th, FOREST, corner=0.06)
        add_text(s, x, ty + Inches(0.2), tw, Inches(0.65),
                 big, size=22, bold=True, color=GOLD, font=HEAD_FONT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        add_text(s, x + Inches(0.2), ty + Inches(0.95), tw - Inches(0.4),
                 Inches(0.4),
                 lbl, size=10.5, color=SAGE, font=BODY_FONT,
                 align=PP_ALIGN.CENTER, line_spacing=1.2)

    # Hangfire schedule table
    sy = Inches(3.55); sw = Inches(7.4); sh = Inches(3.35)
    add_round(s, Inches(0.5), sy, sw, sh, PAPER, corner=0.04, line=DIVIDER)
    add_text(s, Inches(0.7), sy + Inches(0.15), sw - Inches(0.4), Inches(0.4),
             "Hangfire schedule",
             size=14, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
    headers = ["Job", "Cadence"]
    rows = [
        ("RiskEngineService",            "every hour"),
        ("DeveloperScoringService",      "nightly"),
        ("AlertDeliveryService",         "every minute"),
        ("AlertEscalationService",       "every 5 minutes"),
        ("KpiSnapshotRefreshService",    "every 15 minutes"),
        ("InvestorAlertEvaluator",       "every 10 minutes"),
        ("EscrowHealthReportService",    "month-end"),
        ("BenchmarkRefreshService",      "hourly"),
        ("CurrencyRatesRefreshService",  "hourly"),
    ]
    add_table(s, Inches(0.7), sy + Inches(0.6), sw - Inches(0.4), sh - Inches(0.8),
              headers, rows, font_size=10.5,
              col_widths=[Inches(4.6), Inches(2.4)])

    # Right column: forward QA strategy
    rx = Inches(8.1); ry = Inches(3.55); rw = Inches(4.7); rh = Inches(3.35)
    add_round(s, rx, ry, rw, rh, SAGE_SOFT, corner=0.04)
    add_text(s, rx + Inches(0.25), ry + Inches(0.15), rw - Inches(0.5), Inches(0.4),
             "Test strategy across the 4 months",
             size=13, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
    qa_rows = [
        ("Unit (xUnit)",       "per CQRS handler and service"),
        ("Integration",        "per controller against EF Core"),
        ("End-to-end UI",      "Playwright on the critical journeys"),
        ("Load · performance", "k6 / NBomber to P95 sign-off"),
        ("Security · VAPT",    "OWASP ZAP + DESC-authorised VAPT"),
        ("Accessibility",      "axe-core scan per release"),
        ("Data accuracy",      "Great Expectations / dbt-tests"),
        ("Language & RTL",     "native-speaker review per locale"),
    ]
    rly = ry + Inches(0.65)
    for k, v in qa_rows:
        add_text(s, rx + Inches(0.3), rly, Inches(1.85), Inches(0.3),
                 k, size=10.5, bold=True, color=INK, font=BODY_FONT, line_spacing=1.1)
        add_text(s, rx + Inches(1.95), rly, rw - Inches(2.15), Inches(0.6),
                 v, size=9.5, color=INK_SOFT, font=BODY_FONT, line_spacing=1.2)
        rly += Inches(0.32)

    add_footer(s, prs, num, total)


def slide_value_add(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "Vendor commitments", "15 improvements · committed within the 4-month plan")

    commits = [
        ("VC-1",  "MediatR ValidationBehavior",          "Auto-runs every FluentValidation validator."),
        ("VC-2",  "Global IExceptionHandler",            "RFC 7807 ProblemDetails contract."),
        ("VC-3",  "Distributed cache via Redis",         "ICacheableQuery + CachingBehavior."),
        ("VC-4",  "Domain event dispatcher",             "Decouples audit / notify / SignalR."),
        ("VC-5",  "Notification outbox",                 "At-least-once on SMTP/SMS failure."),
        ("VC-6",  "SignalR notification hub",            "Already shipped — bell updates push."),
        ("VC-7",  "AI streaming over SSE",               "IAsyncEnumerable<string> tokens."),
        ("VC-8",  "Polly resilience",                    "AddStandardResilienceHandler() on every HttpClient."),
        ("VC-9",  "Hangfire idempotency",                "Fingerprints prevent double-execution."),
        ("VC-10", "OpenTelemetry",                       "OTLP traces · metrics · logs."),
        ("VC-11", "/healthz + /readyz",                  "Liveness + readiness probes."),
        ("VC-12", "Tighter CORS + CSP + HSTS",           "Replaces AllowAnyOrigin."),
        ("VC-13", "Architecture tests (NetArchTest)",    "Layer-dependency assertions."),
        ("VC-14", "Dockerfiles + CI/CD",                 "Already shipped — multi-stage builds."),
        ("VC-15", "Resource-based authorisation",        "Investor sees own data only."),
    ]
    # Subtitle line, before the card grid, so the bottom of the slide
    # stays free for the footer.
    add_text(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.3),
             "Two items below are already shipped (VC-6 SignalR hub, VC-14 Dockerfiles + CI base); the remaining 13 are committed within the Phase 1 hardening scope.",
             size=10.5, italic=True, color=MUTED, font=BODY_FONT,
             align=PP_ALIGN.LEFT, line_spacing=1.0)

    # 3-column grid of compact cards
    gx = Inches(0.5); gy = Inches(2.05)
    cw = Inches(4.05); ch = Inches(0.92); gap_x = Inches(0.13); gap_y = Inches(0.08)
    cols = 3
    for i, (tag, name, body) in enumerate(commits):
        r = i // cols; c = i % cols
        x = gx + (cw + gap_x) * c
        y = gy + (ch + gap_y) * r
        add_round(s, x, y, cw, ch, PAPER, corner=0.04, line=DIVIDER)
        add_round(s, x + Inches(0.15), y + Inches(0.2), Inches(0.65),
                  Inches(0.5), FOREST, corner=0.2)
        add_text(s, x + Inches(0.15), y + Inches(0.2), Inches(0.65), Inches(0.5),
                 tag, size=10, bold=True, color=GOLD, font=BODY_FONT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        add_text(s, x + Inches(0.9), y + Inches(0.13), cw - Inches(1.0),
                 Inches(0.3),
                 name, size=11, bold=True, color=INK, font=BODY_FONT, line_spacing=1.0)
        add_text(s, x + Inches(0.9), y + Inches(0.43), cw - Inches(1.0),
                 Inches(0.48),
                 body, size=9.5, color=INK_SOFT, font=BODY_FONT, line_spacing=1.2)

    add_footer(s, prs, num, total)


def slide_financials(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "Cost summary", "Indicative · AED · ex-VAT · Option A: Azure UAE North")

    # Top: three big numbers
    tiles = [
        ("AED 2.58 M", "Total one-time CAPEX",       "Infra + Licences + Application Dev"),
        ("AED 629 K",  "Annual OPEX (Year 1)",       "Includes 12-month warranty & support"),
        ("AED 2.01 M", "Application development",     "1,095 person-days across 4 months"),
    ]
    tx = Inches(0.5); ty = Inches(1.85); tw = Inches(4.05); th = Inches(1.6); gap = Inches(0.13)
    for i, (big, lbl, sub) in enumerate(tiles):
        x = tx + (tw + gap) * i
        add_round(s, x, ty, tw, th, FOREST_DEEP, corner=0.05)
        add_text(s, x + Inches(0.2), ty + Inches(0.25), tw - Inches(0.4),
                 Inches(0.6),
                 big, size=28, bold=True, color=GOLD, font=HEAD_FONT,
                 line_spacing=1.0)
        add_text(s, x + Inches(0.2), ty + Inches(0.85), tw - Inches(0.4),
                 Inches(0.35),
                 lbl, size=12, bold=True, color=PAPER, font=BODY_FONT, line_spacing=1.1)
        add_text(s, x + Inches(0.2), ty + Inches(1.17), tw - Inches(0.4),
                 Inches(0.4),
                 sub, size=10, color=SAGE, font=BODY_FONT, line_spacing=1.2)

    # CAPEX top-level split — detail breakdowns follow on the next slides
    headers = ["Cost category", "Amount (AED)", "Type", "Detail"]
    rows = [
        ("Infrastructure CAPEX (Option A)",  "335,500",   "One-time", "see breakdown"),
        ("Software Licences (Year 1)",       "235,200",   "One-time", "see breakdown"),
        ("Application development (M1–M4)",  "2,007,000", "One-time", "split M1–2 / M3–4"),
        ("Total one-time CAPEX",             "2,577,700", "—",        "—"),
        ("Annual OPEX (Option A)",           "629,200",   "Recurring","see breakdown"),
    ]
    add_table(s, Inches(0.5), Inches(3.65), Inches(12.3), Inches(2.85),
              headers, rows, font_size=11,
              col_widths=[Inches(6.0), Inches(2.4), Inches(1.7), Inches(2.2)])

    add_text(s, Inches(0.5), Inches(6.60), Inches(12.3), Inches(0.4),
             "Blended day rates: Architect AED 2,200 · Senior Developer 1,800 · Mid Developer 1,500 · QA / DevOps 1,400.  The next four slides break each category down line-by-line.",
             size=10.5, italic=True, color=MUTED, font=BODY_FONT, line_spacing=1.3)

    add_footer(s, prs, num, total)


def slide_capex_infra(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "Infrastructure CAPEX",
               "One-time setup · costed across three hosting options")

    headers = ["Infrastructure item", "Option A — Azure UAE North", "Option B — On-Premises", "Option C — Hybrid"]
    rows = [
        ("Compute · App Service / VMs (3 hosts)",       "48,000",   "320,000", "220,000"),
        ("SQL Server (managed / licensed)",             "36,000",   "180,000", "120,000"),
        ("Redis cache cluster",                         "18,000",    "85,000",  "18,000"),
        ("Azure AI Search / vector index",              "24,000",    "65,000",  "24,000"),
        ("Microsoft Fabric / OneLake lakehouse",        "42,000",   "200,000",  "42,000"),
        ("Blob Storage / on-prem NAS",                  "12,000",    "55,000",  "35,000"),
        ("WAF + Front Door / on-prem F5",               "30,000",   "140,000",  "95,000"),
        ("CDN (public static assets)",                  "8,000",     "25,000",  "12,000"),
        ("Private endpoints + site-to-site network",    "14,000",    "60,000",  "45,000"),
        ("DR · cross-region replication (≤ 4 h RPO)",   "28,000",   "180,000",  "95,000"),
        ("SSL certificates (wildcard, 3-year)",         "4,500",      "4,500",   "4,500"),
        ("Monitoring · telemetry (Azure Monitor / Seq)","16,000",    "45,000",  "28,000"),
        ("Initial security hardening + pen-test",       "55,000",    "55,000",  "55,000"),
        ("Subtotal (AED)",                              "335,500", "1,414,500", "793,500"),
    ]
    add_table(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(4.9),
              headers, rows, font_size=10,
              col_widths=[Inches(5.3), Inches(2.4), Inches(2.3), Inches(2.3)])

    add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
             "DLD selects the hosting model at contract signature; the rest of the financials adjust accordingly. AI inference and personal data remain UAE-resident on any option.",
             size=10.5, italic=True, color=MUTED, font=BODY_FONT, line_spacing=1.3)

    add_footer(s, prs, num, total)


def slide_capex_licences(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "Software licences · subscriptions",
               "Year-1 stack · open-source components excluded")

    headers = ["Licence / Subscription", "Vendor", "Unit cost (AED)", "Qty", "Annual (AED)"]
    rows = [
        (".NET 9 runtime + ASP.NET Core",       "Microsoft",     "Open-source",          "—", "0"),
        ("Visual Studio Enterprise (dev seats)","Microsoft",     "12,800 / seat / yr",   "5", "64,000"),
        ("Hangfire Pro",                        "HangfireIO",    "7,200 / server / yr",  "2", "14,400"),
        ("Azure OpenAI Service (tokens)",       "Microsoft",     "Consumption-based",    "—", "90,000 est."),
        ("AI fallback model (Anthropic API)",   "Anthropic",     "Consumption-based",    "—", "40,000 est."),
        ("SMTP relay (SendGrid / Postmark)",    "Twilio / Postmark","2,400 / yr",        "1", "2,400"),
        ("SMS gateway (UAE operator)",          "Operator",      "Per-message",          "—", "18,000 est."),
        ("PDF generation (QuestPDF)",           "QuestPDF",      "2,800 / yr",           "1", "2,800"),
        ("Seq log server (cloud)",              "Datalust",      "3,600 / yr",           "1", "3,600"),
        ("Subtotal (Year 1)",                   "—",             "—",                    "—", "235,200"),
    ]
    add_table(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(4.5),
              headers, rows, font_size=10.5,
              col_widths=[Inches(4.2), Inches(2.6), Inches(2.4), Inches(0.8), Inches(2.3)])

    add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4),
             "Token-priced AI usage will be re-tiered at contract based on agreed monthly volume; estimate above sized on internal traffic projections.",
             size=10.5, italic=True, color=MUTED, font=BODY_FONT, line_spacing=1.3)

    add_footer(s, prs, num, total)


def slide_capex_dev_a(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "Application development · Months 1–2",
               "Phase 1 + Phase 2 deliverables · effort and cost")

    headers = ["Deliverable", "Month", "Effort (PD)", "Cost (AED)"]
    rows = [
        ("External Public Portal (CMS · Dashboard · Transactions · Map · Indices)","M1 · P1","120","216,000"),
        ("Bilingual EN/AR · RTL · i18n pipeline",                                   "M1 · P1", "40", "72,000"),
        ("Slice & Dice Analytics engine",                                           "M1 · P1", "60", "108,000"),
        ("Real Estate AI Agent (Orchestrator · RAG · guardrails)",                  "M1 · P1","100","180,000"),
        ("Investor Notifications (email / SMS / in-platform)",                      "M1 · P1", "60","108,000"),
        ("Open Data API · Developer Portal · API-key management",                   "M1 · P1", "50", "90,000"),
        ("Perf hardening · Redis cache · EF tuning",                                "M1 · P1", "30", "54,000"),
        ("Phase 1 QA · UAT support",                                                "M1 · P1", "30", "54,000"),
        ("Phase 1 subtotal",                                                        "M1",    "490","882,000"),
        ("Watchlist v2 · ESG · GRETI sub-index tracker",                            "M2 · P2", "80","144,000"),
        ("Analytics enhancements · saved views · embed",                            "M2 · P2", "40", "72,000"),
        ("Performance hardening · OpenTelemetry · load-test sign-off",              "M2 · P2", "30", "54,000"),
        ("Phase 2 QA · UAT",                                                        "M2 · P2", "20", "36,000"),
        ("Phase 2 subtotal",                                                        "M2",    "170","306,000"),
        ("M1 + M2 cumulative",                                                      "M1–M2", "660","1,188,000"),
    ]
    add_table(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(4.8),
              headers, rows, font_size=10.5,
              col_widths=[Inches(7.4), Inches(1.5), Inches(1.4), Inches(2.0)])

    add_text(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.4),
             "M1 is the heaviest spend — the External Portal and AI Agent stack land in the first month and underpin every later phase.",
             size=10.5, italic=True, color=MUTED, font=BODY_FONT, line_spacing=1.3)

    add_footer(s, prs, num, total)


def slide_capex_dev_b(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "Application development · Months 3–4",
               "Phase 3 + Phase 4 deliverables · effort and cost")

    headers = ["Deliverable", "Month", "Effort (PD)", "Cost (AED)"]
    rows = [
        ("EWRS (risk engine · L1–L4 escalation · dashboard · playbooks)",           "M3 · P3","100","180,000"),
        ("Developer Performance & Rating + Escrow Monitoring",                      "M3 · P3", "80","144,000"),
        ("Beneficial Ownership · Mortgage analytics · Name Validation",             "M3 · P3", "40", "72,000"),
        ("DESC ISR v3 hardening + VAPT (external consultant)",                      "M3 · P3", "30", "90,000"),
        ("Phase 3 QA · UAT · security acceptance",                                  "M3 · P3", "25", "45,000"),
        ("Phase 3 subtotal",                                                        "M3",    "275","531,000"),
        ("Multilingual extension (ZH · RU · UR · FR · HI · DE) + AI multi-language","M4 · P4", "90","162,000"),
        ("Advanced analytics · AI fine-tuning pipeline",                            "M4 · P4", "50", "90,000"),
        ("Phase 4 QA · UAT · language acceptance",                                  "M4 · P4", "20", "36,000"),
        ("Phase 4 subtotal",                                                        "M4",    "160","288,000"),
        ("Total application development (M1–M4)",                                   "M1–M4","1,095","2,007,000"),
    ]
    add_table(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(4.4),
              headers, rows, font_size=10.5,
              col_widths=[Inches(7.4), Inches(1.5), Inches(1.4), Inches(2.0)])

    add_text(s, Inches(0.5), Inches(6.40), Inches(12.3), Inches(0.4),
             "Phase 3 carries the DESC-authorised VAPT engagement (external consultant, AED 90 K embedded in the line item) — clean report required before Phase 3 UAT sign-off.",
             size=10.5, italic=True, color=MUTED, font=BODY_FONT, line_spacing=1.3)

    add_footer(s, prs, num, total)


def slide_opex(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "Annual OPEX", "Year-1 recurring · three hosting options")

    headers = ["OPEX item", "Option A — Azure UAE", "Option B — On-Premises", "Option C — Hybrid"]
    rows = [
        ("Compute · VM · App Service",                "48,000",  "85,000",  "75,000"),
        ("SQL Server (managed annual)",               "36,000",  "24,000",  "28,000"),
        ("Redis · caching annual",                    "18,000",  "12,000",  "18,000"),
        ("AI model API consumption (annual est.)",   "130,000", "130,000", "130,000"),
        ("OneLake · Fabric storage & compute",        "42,000",  "68,000",  "42,000"),
        ("Monitoring & SIEM",                         "16,000",  "28,000",  "22,000"),
        ("SMTP · SMS channels",                       "20,400",  "20,400",  "20,400"),
        ("Software licences (annual renewal)",        "84,800",  "84,800",  "84,800"),
        ("Warranty & support (12-mo post go-live)",  "180,000", "180,000", "180,000"),
        ("Security patching · monthly scans",         "36,000",  "36,000",  "36,000"),
        ("Data quality · reconciliation tooling",     "18,000",  "18,000",  "18,000"),
        ("Annual OPEX (Year 1)",                     "629,200", "686,200", "654,200"),
    ]
    add_table(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(4.7),
              headers, rows, font_size=10,
              col_widths=[Inches(5.3), Inches(2.4), Inches(2.3), Inches(2.3)])

    add_text(s, Inches(0.5), Inches(6.70), Inches(12.3), Inches(0.4),
             "OPEX starts the month after go-live. Warranty & support pricing applies to all three hosting options at parity.",
             size=10.5, italic=True, color=MUTED, font=BODY_FONT, line_spacing=1.3)

    add_footer(s, prs, num, total)


def slide_milestones(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "Payment milestones",
               "Pay on UAT sign-off · 6 gates across 4 months")

    rows = [
        ("M0", "Contract signature",                                                                "10 %", "AED 200,700"),
        ("M1", "Week 1 — mobilisation complete (infrastructure live · team onboard)",               "10 %", "AED 200,700"),
        ("M2", "End of Month 1 — Phase 1 UAT sign-off (External Portal · AI Agent · Open Data)",   "25 %", "AED 501,750"),
        ("M3", "End of Month 2 — Phase 2 UAT sign-off (ESG · GRETI · Analytics v2)",                "15 %", "AED 301,050"),
        ("M4", "End of Month 3 — Phase 3 UAT sign-off (EWRS · Dev Rating · Escrow · VAPT clean)",   "25 %", "AED 501,750"),
        ("M5", "End of Month 4 — Phase 4 UAT sign-off (Multilingual · Advanced Analytics)",         "15 %", "AED 301,050"),
    ]
    by = Inches(1.85)
    for tag, what, pct, amt in rows:
        h = Inches(0.74)
        add_round(s, Inches(0.5), by, Inches(12.3), h, PAPER, corner=0.04, line=DIVIDER)
        add_round(s, Inches(0.7), by + Inches(0.13), Inches(0.85), Inches(0.5),
                  FOREST, corner=0.2)
        add_text(s, Inches(0.7), by + Inches(0.13), Inches(0.85), Inches(0.5),
                 tag, size=13, bold=True, color=GOLD, font=HEAD_FONT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        add_text(s, Inches(1.75), by + Inches(0.22), Inches(8.0), Inches(0.45),
                 what, size=11.5, color=INK, font=BODY_FONT, line_spacing=1.1)
        add_text(s, Inches(10.0), by + Inches(0.22), Inches(0.9), Inches(0.45),
                 pct, size=12, bold=True, color=FOREST, font=BODY_FONT,
                 align=PP_ALIGN.RIGHT, line_spacing=1.0)
        add_text(s, Inches(11.1), by + Inches(0.22), Inches(1.6), Inches(0.45),
                 amt, size=12, bold=True, color=INK, font=BODY_FONT,
                 align=PP_ALIGN.RIGHT, line_spacing=1.0)
        by += h + Inches(0.08)

    add_text(s, Inches(0.5), by + Inches(0.05), Inches(12.3), Inches(0.3),
             "100 % of AED 2,007,000 development CAPEX is conditional on DLD-signed acceptance certificates.",
             size=10, italic=True, color=MUTED, font=BODY_FONT,
             align=PP_ALIGN.CENTER, line_spacing=1.0)

    add_footer(s, prs, num, total)


def slide_findings(prs, num, total):
    s = add_blank(prs)
    add_header(s, prs, "Cross-check findings", "Reports vs. src/ — full reconciliation")

    # Left column: confirmed
    add_round(s, Inches(0.5), Inches(1.85), Inches(6.3), Inches(5.0),
              SAGE_SOFT, corner=0.04)
    add_text(s, Inches(0.75), Inches(2.05), Inches(5.8), Inches(0.4),
             "Confirmed against the codebase",
             size=14, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)
    add_bullets(s, Inches(0.75), Inches(2.55), Inches(5.8), Inches(4.0), [
        "All ~30 Domain entities (Transaction, Project, EscrowAccount, RiskAlert, ScoringWeight, UserAiMemory, SavedAnalyticsView, …) exist at the cited paths.",
        "Every Application/Features/* folder named in the proposal exists with the matching Commands and Queries.",
        "All Razor pages cited (~23 public + 11 admin) and all controllers (20 in WebAPI · 8 in AdminAPI) match the reports.",
        "AuthorizationPolicies publishes internal.read/edit/manage/admin + external.investor exactly as Appendix A claims.",
        "Rate limiter, ApiKey middleware path, JWT + OIDC dual scheme, MFA flow, unsubscribe + erasure — all verified in code.",
        "0 build warnings · 9 Hangfire jobs registered · /healthz/sla aggregates AI · KPI · notifications.",
    ], size=11, color=INK_SOFT)

    # Right column: discrepancies
    add_round(s, Inches(7.0), Inches(1.85), Inches(5.8), Inches(5.0),
              PAPER, corner=0.04, line=DIVIDER)
    add_rect(s, Inches(7.0), Inches(1.85), Inches(0.08), Inches(5.0), GOLD)
    add_text(s, Inches(7.25), Inches(2.05), Inches(5.4), Inches(0.4),
             "Two clarifications to roll into v2",
             size=14, bold=True, color=FOREST, font=HEAD_FONT, line_spacing=1.0)

    add_text(s, Inches(7.25), Inches(2.55), Inches(5.4), Inches(0.35),
             "1.  Target framework",
             size=12.5, bold=True, color=GOLD, font=BODY_FONT, line_spacing=1.0)
    add_text(s, Inches(7.25), Inches(2.92), Inches(5.4), Inches(1.25),
             "Tech Proposal §11 and Compliance Matrix §11.1 state .NET 8. "
             "Every .csproj targets net9.0 today. Update both reports to .NET 9 / "
             "ASP.NET Core 9 for the final submission.",
             size=11, color=INK_SOFT, font=BODY_FONT, line_spacing=1.3)

    # Thin divider so the gap between items reads as intentional
    add_rect(s, Inches(7.25), Inches(4.25), Inches(5.3), Emu(9525), DIVIDER)

    add_text(s, Inches(7.25), Inches(4.45), Inches(5.4), Inches(0.35),
             "2.  Vendor commitment VC-6",
             size=12.5, bold=True, color=GOLD, font=BODY_FONT, line_spacing=1.0)
    add_text(s, Inches(7.25), Inches(4.82), Inches(5.4), Inches(2.0),
             "SignalR notification hub is listed as 'E — vendor commitment' but "
             "NotificationHub.cs and SignalRNotificationBroadcaster.cs are already "
             "wired in WebAPI. Re-classify as C (compliant today) and keep the "
             "broader resilience-stack commitments where they are.",
             size=11, color=INK_SOFT, font=BODY_FONT, line_spacing=1.3)

    add_footer(s, prs, num, total)


def slide_close(prs):
    s = add_blank(prs)
    sw = prs.slide_width; sh = prs.slide_height
    add_rect(s, 0, 0, sw, sh, FOREST_DEEP)
    add_rect(s, 0, 0, Emu(57150), sh, GOLD)

    add_text(s, Inches(0.9), Inches(1.4), Inches(12), Inches(0.6),
             "VENDOR ACCEPTANCE",
             size=12, bold=True, color=GOLD, font=BODY_FONT, line_spacing=1.0)
    add_text(s, Inches(0.9), Inches(1.8), Inches(12), Inches(2.2),
             "Ready to deliver IRETP\non the terms above.",
             size=48, bold=True, color=PAPER, font=HEAD_FONT, line_spacing=1.1)

    add_text(s, Inches(0.9), Inches(4.4), Inches(12), Inches(2.0),
             "Every RFP clause maps to a live module in the repository. "
             "The Compliance Matrix gives a clause-by-clause traceability "
             "statement, the DESC ISR v3 Appendix is the pre-VAPT "
             "self-assessment, and the Financial Proposal is phase-gated "
             "and pay-on-acceptance.\n\n"
             "The two clarifications surfaced in this deck — .NET 9 and "
             "VC-6 already shipped — will be folded into the v2 submission.",
             size=14, color=SAGE_SOFT, font=BODY_FONT, line_spacing=1.45)

    # Footer band
    add_rect(s, 0, sh - Inches(0.8), sw, Inches(0.8), FOREST)
    add_text(s, Inches(0.9), sh - Inches(0.55), Inches(12), Inches(0.3),
             "DLD-IRETP-2026-001  ·  Integrated Real Estate Transparency Platform  ·  Technical Presentation",
             size=11, color=GOLD, font=BODY_FONT, line_spacing=1.0)


# --- Main ------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        slide_cover,
        slide_programme,
        slide_documents,
        slide_verification,
        slide_architecture,
        slide_stack,
        slide_coverage_table,
        slide_portal,
        slide_analytics,
        slide_powerbi_compare,         # NEW — Power BI vs Blazor table
        slide_ai,
        slide_alerts,
        slide_ewrs,
        slide_developer,
        slide_multilingual,
        slide_security,
        slide_hosting,
        slide_data,
        slide_phases,
        slide_responsibility,          # NEW — Responsibility matrix
        slide_quality,
        slide_value_add,
        # Cost slides removed from the client-facing deck. They are emitted to
        # IRETP_Internal_Costs.pptx instead (see build_internal_costs below)
        # for the partnering-company discussion.
        slide_findings,
        slide_close,
    ]
    total = len(builders)
    for i, fn in enumerate(builders, start=1):
        # Title and close slides don't need page numbers
        if fn is slide_cover or fn is slide_close:
            fn(prs)
        else:
            fn(prs, i, total)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}  ({OUT.stat().st_size:,} bytes, {total} slides)")

    build_internal_costs()


def build_internal_costs():
    """Internal cost-only deck — for the partnering-company discussion.
    Not for DLD submission. Hybrid figures here reflect the corrected
    cost model (compute on-prem + Azure lakehouse + site-to-site link).
    """
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        slide_financials,
        slide_capex_infra,
        slide_capex_licences,
        slide_capex_dev_a,
        slide_capex_dev_b,
        slide_opex,
        slide_milestones,
    ]
    total = len(builders)
    for i, fn in enumerate(builders, start=1):
        fn(prs, i, total)

    prs.save(OUT_INTERNAL)
    print(f"Wrote {OUT_INTERNAL}  ({OUT_INTERNAL.stat().st_size:,} bytes, {total} slides)")


if __name__ == "__main__":
    main()
