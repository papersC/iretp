"""Append 3 slides to IRETP_Technical_Presentation.pptx:
  25) Architecture Diagram
  26) CLA (Compliance List Acknowledgement) Table
  27) Hardware Requirements
Re-numbers page footers across all slides to "X / 27".
"""
import sys, io, copy, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

SRC = r"C:\Users\kalmi\IRETP\IRETP_Technical_Presentation.pptx"
OUT = r"C:\Users\kalmi\IRETP\IRETP_Technical_Presentation.pptx"

# ---- DLD palette ----
GREEN_DARK   = RGBColor(0x06, 0x35, 0x25)
GREEN_MAIN   = RGBColor(0x0E, 0x4E, 0x37)
GREEN_MOSS   = RGBColor(0x7C, 0xB9, 0x9F)
GREEN_PANEL  = RGBColor(0xE6, 0xF1, 0xEB)
GOLD         = RGBColor(0xC8, 0xA4, 0x64)
DIVIDER      = RGBColor(0xD7, 0xDE, 0xDA)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK    = RGBColor(0x18, 0x2A, 0x22)
TEXT_MUTED   = RGBColor(0x5A, 0x6C, 0x64)

FONT_HEAD = "Calibri"
FONT_BODY = "Calibri"

# ---- helpers ----
def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()

def set_line(shape, rgb, width_pt=0.75):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is not None:
        set_fill(s, fill)
    else:
        s.fill.background()
    if line is None:
        s.line.fill.background()
    else:
        set_line(s, line, line_w)
    return s

def add_text(slide, x, y, w, h, text, *, size=12, bold=False, italic=False,
             color=TEXT_DARK, font=FONT_BODY, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        f = run.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return tb

def add_page_chrome(slide, title, subtitle, page_no, total):
    # title bar
    add_text(slide, 0.50, 0.45, 12.00, 0.40, title,
             size=20, bold=True, color=GREEN_DARK, font=FONT_HEAD)
    add_text(slide, 0.50, 0.85, 12.00, 0.55, subtitle,
             size=13, italic=True, color=TEXT_MUTED)
    # divider
    add_rect(slide, 0.50, 1.50, 12.30, 0.01, fill=DIVIDER)
    # footer
    add_text(slide, 0.50, 7.14, 8.00, 0.25,
             "IRETP — Technical Presentation",
             size=9, color=TEXT_MUTED)
    add_text(slide, 12.03, 7.14, 0.90, 0.25,
             f"{page_no} / {total}",
             size=9, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)

# =============================================================
# SLIDE A — Architecture Diagram
# =============================================================
def build_architecture(slide, page_no, total):
    add_page_chrome(slide,
        "ARCHITECTURE DIAGRAM",
        "Runtime topology · request flow · UAE-resident perimeter",
        page_no, total)

    # ---- Tier-1: Clients ----
    y = 1.80
    add_text(slide, 0.50, y, 3.00, 0.28, "CLIENTS",
             size=9, bold=True, color=TEXT_MUTED, font=FONT_HEAD)
    add_rect(slide, 0.50, y+0.30, 3.00, 0.85, fill=WHITE, line=DIVIDER)
    add_text(slide, 0.50, y+0.32, 3.00, 0.30, "Investors · DLD Staff · Developers",
             size=11, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, 0.50, y+0.62, 3.00, 0.50, "Browser (Blazor)  ·  Mobile Web\nPublic API consumers",
             size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # ---- Tier-2: Edge (Front Door + WAF) ----
    add_text(slide, 3.85, y, 2.50, 0.28, "EDGE & SECURITY",
             size=9, bold=True, color=TEXT_MUTED, font=FONT_HEAD)
    edge = add_rect(slide, 3.85, y+0.30, 2.50, 0.85, fill=GREEN_PANEL, line=GREEN_MOSS)
    add_text(slide, 3.85, y+0.32, 2.50, 0.30, "Front Door + WAF",
             size=11, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, 3.85, y+0.62, 2.50, 0.50, "OWASP rules · DDoS L7\nTLS 1.2+  ·  geo policy",
             size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # ---- Tier-3: Identity ----
    add_text(slide, 6.70, y, 2.50, 0.28, "IDENTITY",
             size=9, bold=True, color=TEXT_MUTED, font=FONT_HEAD)
    add_rect(slide, 6.70, y+0.30, 2.50, 0.85, fill=WHITE, line=DIVIDER)
    add_text(slide, 6.70, y+0.32, 2.50, 0.30, "OIDC + JWT",
             size=11, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, 6.70, y+0.62, 2.50, 0.50, "ASP.NET Identity · TOTP MFA\nUAE PASS-ready",
             size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # ---- Tier-4: Observability ----
    add_text(slide, 9.55, y, 3.30, 0.28, "OBSERVABILITY",
             size=9, bold=True, color=TEXT_MUTED, font=FONT_HEAD)
    add_rect(slide, 9.55, y+0.30, 3.30, 0.85, fill=WHITE, line=DIVIDER)
    add_text(slide, 9.55, y+0.32, 3.30, 0.30, "App Insights · Log Analytics",
             size=11, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, 9.55, y+0.62, 3.30, 0.50, "SLO health · /healthz/sla\nAudit trail (append-only)",
             size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # ---- Tier-5: Application hosts (3 boxes) ----
    yh = 3.10
    add_text(slide, 0.50, yh-0.20, 12.30, 0.25, "APPLICATION HOSTS  ·  ASP.NET Core (.NET 9)",
             size=10, bold=True, color=GREEN_DARK, font=FONT_HEAD)

    host_specs = [
        ("IRETP.Web",     "Blazor Server portal",     "Public + Admin Razor pages",  0.50),
        ("IRETP.WebAPI",  "Public + Investor API",    "JWT · rate-limited · /v1/open-data/*",  4.60),
        ("IRETP.AdminAPI","Internal DLD operations",  "EWRS · Escrow · Rating · private network",  8.70),
    ]
    for name, sub, det, x in host_specs:
        # gold top accent
        add_rect(slide, x, yh+0.05, 4.10, 0.07, fill=GOLD)
        # main box
        add_rect(slide, x, yh+0.05, 4.10, 1.05, fill=GREEN_MAIN)
        add_text(slide, x+0.20, yh+0.18, 3.70, 0.30, name,
                 size=13, bold=True, color=WHITE, font=FONT_HEAD)
        add_text(slide, x+0.20, yh+0.50, 3.70, 0.28, sub,
                 size=10, color=WHITE)
        add_text(slide, x+0.20, yh+0.78, 3.70, 0.28, det,
                 size=9, italic=True, color=GREEN_PANEL)

    # connector lines from edge into hosts (small visual)
    # we'll just rely on layered placement; explicit arrows can clutter

    # ---- Tier-6: Application core band (CQRS) ----
    yc = 4.35
    band = add_rect(slide, 0.50, yc, 12.30, 0.85, fill=GREEN_PANEL)
    add_text(slide, 0.70, yc+0.08, 11.90, 0.28,
             "IRETP.Application  ·  CQRS via MediatR  ·  Feature folders (AIAgent · Alerts · Analytics · EWRS · Escrow · Map · Transactions · …)",
             size=10, bold=True, color=GREEN_DARK)
    add_text(slide, 0.70, yc+0.38, 11.90, 0.42,
             "IRETP.Domain  ·  DDD entities + enums  ·  Transaction · Project · EscrowAccount · RiskAlert · ScoringWeight · BeneficialOwner",
             size=10, color=TEXT_DARK)

    # ---- Tier-7: Infrastructure / data services ----
    yi = 5.35
    add_text(slide, 0.50, yi, 12.30, 0.25,
             "INFRASTRUCTURE  ·  IRETP.Infrastructure",
             size=10, bold=True, color=GREEN_DARK, font=FONT_HEAD)

    infra = [
        ("SQL Server",   "OLTP store · EF Core 9", 0.50),
        ("OneLake",      "Lakehouse · open data",  2.55),
        ("Hangfire",     "Background jobs · SLA",  4.60),
        ("AI Orchestr.", "UAE-resident LLM tiers", 6.65),
        ("Key Vault",    "Secrets · MSI · rotation",8.70),
        ("Blob (GRS)",   "CMS · exports · backups",10.75),
    ]
    for name, sub, x in infra:
        add_rect(slide, x, yi+0.30, 2.00, 0.80, fill=WHITE, line=GREEN_MOSS)
        add_rect(slide, x, yi+0.30, 2.00, 0.06, fill=GOLD)
        add_text(slide, x, yi+0.40, 2.00, 0.28, name,
                 size=10, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER, font=FONT_HEAD)
        add_text(slide, x, yi+0.66, 2.00, 0.42, sub,
                 size=8, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # ---- Bottom callout ----
    add_rect(slide, 0.50, 6.40, 12.30, 0.62, fill=GREEN_DARK)
    add_text(slide, 0.70, 6.46, 11.90, 0.25,
             "UAE-RESIDENT PERIMETER  ·  TLS 1.2+ end-to-end  ·  Managed Identity everywhere  ·  No cross-border data flow",
             size=10, bold=True, color=GOLD, font=FONT_HEAD)
    add_text(slide, 0.70, 6.72, 11.90, 0.25,
             "Healthz · /healthz/live · /healthz/ready · /healthz/sla   |   RTO < 4 h · RPO < 1 h   |   DESC ISR v3 controls mapped",
             size=9, color=WHITE)

# =============================================================
# SLIDE B — CLA (Compliance List Acknowledgement) Table
# =============================================================
def build_cla_table(slide, page_no, total):
    add_page_chrome(slide,
        "COMPLIANCE LIST ACKNOWLEDGEMENT (CLA)",
        "Per-section compliance status against RFP DLD-IRETP-2026-001",
        page_no, total)

    # Legend pills
    pills = [("FC", "Full Compliance",    GREEN_MAIN, WHITE),
             ("PC", "Partial Compliance", GOLD,       GREEN_DARK),
             ("NC", "Non-Compliance",     RGBColor(0xB0,0x42,0x42), WHITE)]
    px = 0.50
    add_text(slide, 0.50, 1.55, 1.20, 0.28, "Legend:",
             size=9, bold=True, color=TEXT_MUTED, font=FONT_HEAD)
    px = 1.20
    for code, label, bg, fg in pills:
        add_rect(slide, px, 1.57, 0.34, 0.22, fill=bg)
        add_text(slide, px, 1.57, 0.34, 0.22, code,
                 size=9, bold=True, color=fg, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, px+0.40, 1.55, 1.55, 0.28, label,
                 size=9, color=TEXT_DARK)
        px += 2.05

    # ---- Table ----
    rows = [
        ("§3",  "External Public Portal",          "12", "FC", "Blazor Server pages live · CMS · bilingual · KPI dashboard"),
        ("§4",  "Slice-and-Dice Analytics",        "9",  "FC", "Saved views (cap 12) · 4 FR-004 charts · Zone Compare (≤5)"),
        ("§5",  "Real Estate AI Agent",            "11", "FC", "UAE-resident · multi-model tier routing · advisory guardrail · 14-case harness"),
        ("§6",  "Investor Notifications",          "8",  "FC", "Email/SMS/in-platform · SLA probe · HMAC unsubscribe · List-Unsubscribe"),
        ("§7",  "Multilingual EN/AR + Currency",   "6",  "FC", "RTL parity · Phase-4 ZH/RU/UR/FR/HI/DE roadmap"),
        ("§8",  "Early Warning Risk System",       "10", "FC", "10 indicators · 4-level auto-escalation · threshold-driven · admin override"),
        ("§9",  "Developer Rating & Escrow",       "8",  "FC", "Public scorecards · §8.4 escrow monitor · weighted composite"),
        ("§10", "Security & RBAC (DESC ISR v3)",   "14", "FC", "MFA TOTP · policy registry · audit immutability · SLO health"),
        ("§11", "Hosting & Disaster Recovery",     "9",  "FC", "Three options · RTO < 4 h · RPO < 1 h · IaC (Bicep) · CI sec-scan"),
        ("§13", "Code Quality & Documentation",    "7",  "FC", "88/88 xUnit pass · 0 warn / 0 err · ARCH · API · RUNBOOK · DR · DICT"),
        ("§20", "ESG & Open Data",                 "5",  "FC", "ESG heatmap · /v1/open-data/* · OpenAPI 3.0 · API keys"),
    ]

    table_top = 1.95
    table_left = 0.50
    table_w = 12.30
    header_h = 0.38
    row_h = 0.30

    col_w = [0.95, 3.30, 1.10, 1.20, table_w - 0.95 - 3.30 - 1.10 - 1.20]
    headers = ["RFP §", "Section Title", "# Reqs", "Status", "Evidence / Notes"]

    # header
    add_rect(slide, table_left, table_top, table_w, header_h, fill=GREEN_DARK)
    cx = table_left
    for i, hd in enumerate(headers):
        align = PP_ALIGN.CENTER if i in (0,2,3) else PP_ALIGN.LEFT
        add_text(slide, cx+0.10, table_top+0.06, col_w[i]-0.20, header_h-0.10, hd,
                 size=11, bold=True, color=GOLD, font=FONT_HEAD,
                 align=align, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[i]

    # rows
    ry = table_top + header_h
    for r_idx, (sec, title, n, status, notes) in enumerate(rows):
        zebra = GREEN_PANEL if r_idx % 2 == 0 else WHITE
        add_rect(slide, table_left, ry, table_w, row_h, fill=zebra)
        cx = table_left
        cells = [sec, title, n, status, notes]
        for i, val in enumerate(cells):
            align = PP_ALIGN.CENTER if i in (0,2,3) else PP_ALIGN.LEFT
            if i == 3:
                # status pill
                pill_w = 0.60
                pill_x = cx + (col_w[i] - pill_w) / 2
                pill_h = 0.22
                pill_y = ry + (row_h - pill_h) / 2
                bg, fg = GREEN_MAIN, WHITE
                if val == "PC": bg, fg = GOLD, GREEN_DARK
                elif val == "NC": bg, fg = RGBColor(0xB0,0x42,0x42), WHITE
                add_rect(slide, pill_x, pill_y, pill_w, pill_h, fill=bg)
                add_text(slide, pill_x, pill_y, pill_w, pill_h, val,
                         size=10, bold=True, color=fg,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
            else:
                size = 10 if i != 4 else 9
                bold = (i == 0 or i == 1)
                add_text(slide, cx+0.10, ry, col_w[i]-0.20, row_h, str(val),
                         size=size, bold=bold,
                         color=GREEN_DARK if i in (0,1) else TEXT_DARK,
                         align=align, anchor=MSO_ANCHOR.MIDDLE,
                         font=FONT_HEAD if i in (0,1) else FONT_BODY)
            cx += col_w[i]
        ry += row_h

    # Total bar
    total_reqs = sum(int(r[2]) for r in rows)
    add_rect(slide, table_left, ry, table_w, 0.38, fill=GREEN_DARK)
    add_text(slide, table_left+0.10, ry, table_w-0.20, 0.38,
             f"TOTAL  ·  {len(rows)} sections  ·  {total_reqs} requirements  ·  "
             f"{len(rows)}/{len(rows)} Full Compliance  ·  100 %",
             size=11, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)

    # Bottom callout — must end before footer (y=7.14)
    cy = ry + 0.48
    add_rect(slide, 0.50, cy, 12.30, 0.55, fill=GREEN_PANEL)
    add_text(slide, 0.70, cy+0.04, 11.90, 0.22,
             "Verification — 88/88 xUnit tests pass · 0 build warnings · 0 build errors · cross-check vs src/ reconciled",
             size=10, bold=True, color=GREEN_DARK, font=FONT_HEAD)
    add_text(slide, 0.70, cy+0.27, 11.90, 0.26,
             "Full traceability in Compliance Matrix (Vol 2) — every line maps to a Domain entity, Application handler and Razor page or controller in src/.",
             size=9, italic=True, color=TEXT_MUTED)

# =============================================================
# SLIDE C — Hardware Requirements
# =============================================================
def build_hardware(slide, page_no, total):
    add_page_chrome(slide,
        "HARDWARE REQUIREMENTS",
        "Production sizing · per hosting option · UAE-resident",
        page_no, total)

    # ---- Application tier table ----
    y_app = 1.75
    add_text(slide, 0.50, y_app, 12.30, 0.30,
             "APPLICATION & DATA TIER (steady-state production)",
             size=11, bold=True, color=GREEN_DARK, font=FONT_HEAD)

    headers = ["Component", "Role", "vCPU", "RAM", "Storage", "Instances", "Azure SKU equivalent"]
    rows = [
        ("IRETP.Web",        "Blazor Server portal",        "4",  "16 GB", "64 GB SSD",    "2 (HA)", "D4s_v5"),
        ("IRETP.WebAPI",     "Public + Investor API",       "4",  "16 GB", "64 GB SSD",    "2 (HA)", "D4s_v5"),
        ("IRETP.AdminAPI",   "Internal DLD operations",     "2",  "8 GB",  "64 GB SSD",    "2 (HA)", "D2s_v5"),
        ("Hangfire worker",  "EWRS · alerts · KPI refresh", "4",  "16 GB", "128 GB SSD",   "2 (HA)", "D4s_v5"),
        ("SQL Server",       "Primary OLTP",                "8",  "64 GB", "1 TB Premium SSD", "1 + replica", "MI BC 8 vCore"),
        ("OneLake / Blob",   "Lakehouse + exports + CMS",   "—",  "—",     "5 TB GRS",     "—",      "Storage v2 (GRS)"),
        ("Key Vault",        "Secrets · MSI · rotation",    "—",  "—",     "—",            "1",      "Premium (HSM)"),
        ("Front Door + WAF", "Edge · OWASP · DDoS",         "—",  "—",     "—",            "Global", "Premium"),
    ]
    table_top = y_app + 0.35
    table_left = 0.50
    table_w = 12.30
    header_h = 0.36
    row_h = 0.32
    col_w = [1.85, 2.90, 0.70, 0.90, 1.85, 1.30, table_w - 1.85 - 2.90 - 0.70 - 0.90 - 1.85 - 1.30]

    add_rect(slide, table_left, table_top, table_w, header_h, fill=GREEN_DARK)
    cx = table_left
    for i, hd in enumerate(headers):
        align = PP_ALIGN.CENTER if i in (2,3,4,5) else PP_ALIGN.LEFT
        add_text(slide, cx+0.08, table_top+0.04, col_w[i]-0.16, header_h-0.08, hd,
                 size=10, bold=True, color=GOLD, font=FONT_HEAD,
                 align=align, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[i]

    ry = table_top + header_h
    for r_idx, row in enumerate(rows):
        zebra = GREEN_PANEL if r_idx % 2 == 0 else WHITE
        add_rect(slide, table_left, ry, table_w, row_h, fill=zebra)
        cx = table_left
        for i, val in enumerate(row):
            align = PP_ALIGN.CENTER if i in (2,3,4,5) else PP_ALIGN.LEFT
            bold = (i == 0)
            add_text(slide, cx+0.08, ry, col_w[i]-0.16, row_h, str(val),
                     size=9, bold=bold,
                     color=GREEN_DARK if i == 0 else TEXT_DARK,
                     align=align, anchor=MSO_ANCHOR.MIDDLE,
                     font=FONT_HEAD if i == 0 else FONT_BODY)
            cx += col_w[i]
        ry += row_h

    # ---- Three option footprints (mini cards) ----
    y_opt = ry + 0.20
    add_text(slide, 0.50, y_opt, 12.30, 0.28,
             "TOTAL FOOTPRINT BY HOSTING OPTION",
             size=11, bold=True, color=GREEN_DARK, font=FONT_HEAD)

    opt_cards = [
        ("Option A — Azure UAE North",
         "Managed PaaS · MSI · GRS backup",
         "App: P1v3 (~12 vCPU / 48 GB)",
         "Data: MI BC 8 vCore  ·  5 TB GRS",
         "RPO ≤ 1 h  ·  RTO < 4 h"),
        ("Option B — On-Prem (DLD DC)",
         "Bare-metal VMs · DLD-managed",
         "App: 4 hosts × 8 vCPU / 32 GB",
         "Data: SQL 2022 Cluster  ·  10 TB SAN",
         "RPO ≤ 1 h  ·  RTO < 4 h"),
        ("Option C — Hybrid",
         "Compute on-prem · lakehouse Azure",
         "App: 3 hosts × 8 vCPU / 32 GB",
         "Data: SQL on-prem + OneLake (UAE)",
         "RPO ≤ 1 h  ·  RTO < 4 h"),
    ]
    cy = y_opt + 0.32
    card_h = 1.00
    for i, (title, sub, l1, l2, l3) in enumerate(opt_cards):
        x = 0.50 + i * 4.10
        # gold accent
        add_rect(slide, x, cy, 4.00, 0.06, fill=GOLD)
        # body
        add_rect(slide, x, cy, 4.00, card_h, fill=WHITE, line=DIVIDER)
        add_text(slide, x+0.15, cy+0.08, 3.70, 0.24, title,
                 size=10, bold=True, color=GREEN_DARK, font=FONT_HEAD)
        add_text(slide, x+0.15, cy+0.30, 3.70, 0.20, sub,
                 size=8, italic=True, color=TEXT_MUTED)
        add_text(slide, x+0.15, cy+0.50, 3.70, 0.18, l1,
                 size=8, color=TEXT_DARK)
        add_text(slide, x+0.15, cy+0.66, 3.70, 0.18, l2,
                 size=8, color=TEXT_DARK)
        add_text(slide, x+0.15, cy+0.82, 3.70, 0.18, l3,
                 size=8, bold=True, color=GREEN_MAIN)

    # ---- Bottom mandatory controls callout (must end before footer y=7.14) ----
    by = cy + card_h + 0.12
    add_rect(slide, 0.50, by, 12.30, 0.46, fill=GREEN_DARK)
    add_text(slide, 0.70, by+0.03, 11.90, 0.20,
             "MANDATORY — every option",
             size=9, bold=True, color=GOLD, font=FONT_HEAD)
    add_text(slide, 0.70, by+0.22, 11.90, 0.22,
             "TLS 1.2+ end-to-end  ·  segregated networks (Web/API/Admin)  ·  daily encrypted backups  ·  HSM-backed keys  ·  WAF + DDoS L7  ·  UAE-resident",
             size=9, color=WHITE)

# =============================================================
# MAIN
# =============================================================
def main():
    prs = Presentation(SRC)
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    # Find a blank layout
    for L in prs.slide_layouts:
        if L.name.lower() == "blank":
            blank_layout = L
            break

    existing = len(prs.slides)            # 24
    new_total = existing + 3              # 27

    # Build new slides
    s1 = prs.slides.add_slide(blank_layout)
    build_architecture(s1, existing + 1, new_total)
    s2 = prs.slides.add_slide(blank_layout)
    build_cla_table(s2, existing + 2, new_total)
    s3 = prs.slides.add_slide(blank_layout)
    build_hardware(s3, existing + 3, new_total)

    # Renumber footers on existing slides: replace "X / 24" or "X/24" with "X / 27"
    for idx, slide in enumerate(prs.slides, start=1):
        if idx > existing:
            continue  # new slides are already correct
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = shape.text_frame.text
            if "/" not in txt:
                continue
            # match patterns "N / 24" or "N/24"
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    rt = r.text
                    if "/ 24" in rt or "/24" in rt:
                        # try to retain N
                        r.text = rt.replace("/ 24", f"/ {new_total}").replace("/24", f"/{new_total}")

    # Save
    try:
        prs.save(OUT)
        print(f"OK  saved -> {OUT}  (slides: {existing} -> {new_total})")
    except PermissionError as e:
        alt = OUT.replace(".pptx", "_updated.pptx")
        prs.save(alt)
        print(f"LOCKED  original is open in PowerPoint.  Saved fallback -> {alt}")
        print(f"        Close PowerPoint, delete the lock file, then move {alt} over the original.")

if __name__ == "__main__":
    main()
