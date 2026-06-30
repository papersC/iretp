# -*- coding: utf-8 -*-
"""
IRETP / DLD-IRETP-2026-001 — Technical Proposal Presentation builder.
Premium DLD-branded deck: nano-banana hero images, framed real screenshots,
vector architecture + timeline diagrams, speaker notes. 16:9.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from PIL import Image, ImageDraw

ROOT = r"C:\Users\kalmi\IRETP"
ASSETS = os.path.join(ROOT, "assets")
SCREENS_DIR = os.path.join(ROOT, "proposal_screens")
FRAMED = os.path.join(ASSETS, "framed")
os.makedirs(FRAMED, exist_ok=True)
OUT = os.path.join(ROOT, "IRETP_DLD_Technical_Presentation.pptx")

# ---- Brand palette -------------------------------------------------------
INK    = "10271C"   # near-black green (body text on light)
GREEN  = "0A6E3D"   # DLD primary green
GREEN2 = "0E5A35"   # secondary green
BRIGHT = "1AA45B"   # bright accent green
DARK   = "06281A"   # dark slide base / scrim
GOLD   = "C6A24A"   # premium gold accent
GOLD2  = "B8932E"
MIST   = "ECF3EE"   # light green-tint card
CARD   = "F4F8F5"
NUMTINT= "E4EEE7"   # big translucent number
MUTE   = "5E6F66"   # muted captions
LINE   = "D9E5DE"   # hairline
WHITE  = "FFFFFF"
PALE   = "CFE3D6"   # pale green text on dark

# ---- Fonts (Windows-native, verified installed) --------------------------
HEAVY = "Segoe UI Black"
SEMI  = "Segoe UI Semibold"
BODY  = "Segoe UI"
LIGHT = "Segoe UI Light"

EMU_IN = 914400
SW, SH = 13.333, 7.5

# ---- Screenshot rounding -------------------------------------------------
def make_rounded(src, dst, frac=0.022):
    im = Image.open(src).convert("RGBA")
    w, h = im.size
    rad = int(w * frac)
    mask = Image.new("L", (w, h), 0)
    ImageDraw.Draw(mask).rounded_rectangle([0, 0, w, h], radius=rad, fill=255)
    im.putalpha(mask)
    im.save(dst)

SCREENS = {
    "dashboard": "01_dashboard.png", "map": "02_map.png", "transactions": "03_transactions.png",
    "ai": "04_ai_agent.png", "ewrs": "05_ewrs.png", "scorecards": "06_scorecards.png",
    "bo": "07_beneficial_ownership.png", "mortgage": "08_mortgage.png", "audit": "09_audit_logs.png",
}
for k, fn in SCREENS.items():
    s = os.path.join(SCREENS_DIR, fn)
    d = os.path.join(FRAMED, k + ".png")
    if os.path.exists(s):
        make_rounded(s, d)
SHOT_AR = 2400 / 1380.0  # width / height

# ---- Presentation --------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, hexc):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor.from_string(hexc)

def _clear_fills(spPr):
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        for e in spPr.findall(qn(tag)):
            spPr.remove(e)

def _insert_fill(shape, xml):
    spPr = shape._element.spPr
    _clear_fills(spPr)
    el = parse_xml(xml)
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        ln.addprevious(el)
    else:
        spPr.append(el)

def solid_alpha(shape, hexc, opacity):
    _insert_fill(shape, f'<a:solidFill {nsdecls("a")}><a:srgbClr val="{hexc}"><a:alpha val="{int(opacity*1000)}"/></a:srgbClr></a:solidFill>')

def gradient(shape, stops, ang=5400000):
    gs = "".join(f'<a:gs pos="{int(p*1000)}"><a:srgbClr val="{h}"><a:alpha val="{int(o*1000)}"/></a:srgbClr></a:gs>' for p, h, o in stops)
    _insert_fill(shape, f'<a:gradFill {nsdecls("a")} rotWithShape="1"><a:gsLst>{gs}</a:gsLst><a:lin ang="{ang}" scaled="1"/></a:gradFill>')

def shadow(shape, blur=9, dist=4, dir_deg=90, alpha=34, color="06281A"):
    spPr = shape._element.spPr
    for e in spPr.findall(qn("a:effectLst")):
        spPr.remove(e)
    spPr.append(parse_xml(
        f'<a:effectLst {nsdecls("a")}><a:outerShdw blurRad="{int(blur*12700)}" dist="{int(dist*12700)}" '
        f'dir="{int(dir_deg*60000)}" rotWithShape="0"><a:srgbClr val="{color}"><a:alpha val="{int(alpha*1000)}"/>'
        f'</a:srgbClr></a:outerShdw></a:effectLst>'))

def rect(s, l, t, w, h, fill=None, line=None, lw=1.0, rad=0.0):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rad > 0 else MSO_SHAPE.RECTANGLE
    o = s.shapes.add_shape(shp, Inches(l), Inches(t), Inches(w), Inches(h))
    if rad > 0:
        try: o.adjustments[0] = rad
        except Exception: pass
    o.shadow.inherit = False
    if fill is None:
        o.fill.background()
    else:
        o.fill.solid(); o.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        o.line.fill.background()
    else:
        o.line.color.rgb = RGBColor.from_string(line); o.line.width = Pt(lw)
    return o

def line_h(s, l, t, w, color=LINE, weight=1.0):
    o = s.shapes.add_connector(2, Inches(l), Inches(t), Inches(l + w), Inches(t))
    o.line.color.rgb = RGBColor.from_string(color); o.line.width = Pt(weight)
    o.shadow.inherit = False
    return o

def pic(s, path, l, t, w, h=None, shadowed=False):
    kw = dict(width=Inches(w))
    if h is not None:
        kw["height"] = Inches(h)
    p = s.shapes.add_picture(path, Inches(l), Inches(t), **kw)
    if shadowed:
        shadow(p, blur=11, dist=5, alpha=30)
    return p

def tb(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    return box, tf

def para(tf, text, size, color, font=BODY, bold=False, first=False, align=PP_ALIGN.LEFT,
         before=0, after=0, line=1.0, spacing=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if before: p.space_before = Pt(before)
    p.space_after = Pt(after)
    p.line_spacing = line
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.name = font; f.bold = bold
    f.color.rgb = RGBColor.from_string(color)
    if spacing is not None:
        rPr = r._r.get_or_add_rPr(); rPr.set("spc", str(int(spacing * 100)))
    return p, r

def runs(p, parts):
    """parts: list of (text, size, color, font, bold)"""
    for (text, size, color, font, bold) in parts:
        r = p.add_run(); r.text = text
        f = r.font; f.size = Pt(size); f.name = font; f.bold = bold
        f.color.rgb = RGBColor.from_string(color)

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

def node(s, l, t, size=0.12, color=BRIGHT):
    rect(s, l, t, size, size, fill=color, rad=0.18)

PAGENO = [0]
def footer(s, dark=False):
    PAGENO[0] += 1
    c = PALE if dark else MUTE
    ln = "FFFFFF" if dark else LINE
    line_h(s, 0.92, 7.06, 11.49, color=ln, weight=0.75)
    _, tf = tb(s, 0.92, 7.12, 8, 0.3)
    para(tf, "IRETP  ·  Integrated Real Estate Transparency Platform  ·  DLD-IRETP-2026-001",
         8, c, font=BODY, first=True)
    _, tf2 = tb(s, 11.0, 7.12, 1.41, 0.3)
    para(tf2, f"{PAGENO[0]:02d}", 8, c, font=SEMI, first=True, align=PP_ALIGN.RIGHT)

def header(s, eyebrow, title, num):
    node(s, 0.92, 0.66)
    _, tf = tb(s, 1.16, 0.55, 9.5, 0.35)
    para(tf, eyebrow.upper(), 11, GREEN, font=SEMI, first=True, spacing=2.2)
    _, tf2 = tb(s, 0.9, 0.92, 9.8, 0.85)
    para(tf2, title, 29, INK, font=SEMI, first=True)
    _, tnf = tb(s, 11.1, 0.34, 1.5, 1.0)
    para(tnf, num, 50, NUMTINT, font=HEAVY, first=True, align=PP_ALIGN.RIGHT)

def chip(s, l, t, w, h, text, fill, txtcolor, size=10.5, font=SEMI):
    rect(s, l, t, w, h, fill=fill, rad=0.5)
    _, tf = tb(s, l, t, w, h, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, text, size, txtcolor, font=font, first=True, align=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 1 — COVER
# =========================================================================
def s_cover():
    s = slide()
    pic(s, os.path.join(ASSETS, "cover_hero.png"), 0, 0, SW, SH)
    full = rect(s, 0, 0, SW, SH); gradient(full, [(0, DARK, 18), (60, DARK, 30), (100, DARK, 78)])
    left = rect(s, 0, 0, 8.6, SH); gradient(left, [(0, DARK, 80), (100, DARK, 0)], ang=0)
    # brandmark
    _, tf = tb(s, 0.92, 0.62, 6, 0.9)
    para(tf, "IRETP", 22, WHITE, font=HEAVY, first=True, spacing=1.0)
    para(tf, "INTEGRATED REAL ESTATE TRANSPARENCY PLATFORM", 9, PALE, font=SEMI, before=2, spacing=2.0)
    # reference-build pill (top-right)
    chip(s, 9.5, 0.7, 3.0, 0.42, "WORKING REFERENCE BUILD  ·  ~90% RFP", "FFFFFF", GREEN, size=9, font=SEMI)
    solid_alpha(s.shapes[-2], WHITE, 88)
    # main block
    _, tf2 = tb(s, 0.92, 3.55, 10.4, 3.0)
    para(tf2, "TECHNICAL PROPOSAL PRESENTATION   ·   RFP v1.3", 12.5, GOLD, font=SEMI, first=True, spacing=2.4)
    para(tf2, "Integrated Real Estate", 41, WHITE, font=SEMI, before=10, line=1.02)
    para(tf2, "Transparency Platform", 41, WHITE, font=SEMI, line=1.02)
    para(tf2, "A unified transparency, risk & intelligence platform for Dubai’s real-estate market.",
         15, PALE, font=BODY, before=12, line=1.15)
    # tender + prepared for
    line_h(s, 0.96, 6.62, 5.4, color=GOLD, weight=1.4)
    _, tf3 = tb(s, 0.92, 6.72, 11.5, 0.6)
    runs(para(tf3, "", 12, WHITE, first=True)[0], [
        ("Tender DLD-IRETP-2026-001", 12, WHITE, SEMI, False),
        ("    Prepared for the Dubai Land Department, Government of Dubai", 12, PALE, BODY, False),
    ])
    _, tf4 = tb(s, 0.92, 7.04, 11.5, 0.4)
    runs(para(tf4, "", 11, PALE, first=True)[0], [
        ("[Your Company Name]", 11, GOLD, SEMI, False),
        ("   ·   [Presenter Name, Title]   ·   22 June 2026", 11, PALE, BODY, False),
    ])
    notes(s, "Open warmly. One line: 'Thank you for shortlisting us. In the next 30 minutes we'll show you not a "
              "promise, but a working platform.' Stress the hook: IRETP already exists as a reference build with ~90% "
              "of the RFP demonstrated. This de-risks the whole programme. Then walk the agenda.")

# =========================================================================
# SLIDE 2 — AGENDA
# =========================================================================
def s_agenda():
    s = slide(); bg(s, WHITE)
    header(s, "Agenda", "What we’ll cover today", "")
    _, tf = tb(s, 0.9, 1.62, 11.5, 0.4)
    para(tf, "Eleven sections · roughly 30 minutes · followed by discussion", 13, MUTE, font=BODY, first=True)
    items = [
        ("01", "Dubai Land Department", "Context & mandate"),
        ("02", "Project Summary", "What IRETP is, in one view"),
        ("03", "Project Requirements", "One platform, four audiences"),
        ("04", "Solution Architecture", "Clean Architecture on Microsoft Fabric"),
        ("05", "Delivery Methodology", "How we build and assure"),
        ("06", "Functional Requirements", "What the platform does"),
        ("07", "The Platform — Live Screens", "AI · Maps · a running product"),
        ("08", "Non-Functional Requirements", "Performance, security, residency"),
        ("09", "Implementation Timeline", "3 months to go-live"),
        ("10", "Documentation Provided", "A complete knowledge pack"),
        ("11", "Project Deliverables", "What DLD receives"),
    ]
    cols, x0, y0, cw, ch, gx, gy = 2, 0.9, 2.25, 5.72, 0.62, 0.1, 0.16
    for i, (n, t, sub) in enumerate(items):
        c, r = i % cols, i // cols
        x, y = x0 + c * (cw + gx), y0 + r * (ch + gy)
        rect(s, x, y, cw, ch, fill=CARD, rad=0.10)
        rect(s, x, y, 0.07, ch, fill=BRIGHT)
        _, nt = tb(s, x + 0.22, y, 0.7, ch, anchor=MSO_ANCHOR.MIDDLE)
        para(nt, n, 17, GREEN, font=HEAVY, first=True)
        _, it = tb(s, x + 0.95, y + 0.08, cw - 1.05, ch - 0.1, anchor=MSO_ANCHOR.MIDDLE)
        para(it, t, 13, INK, font=SEMI, first=True, line=1.0)
        para(it, sub, 9.5, MUTE, font=BODY, before=1, line=1.0)
    footer(s)
    notes(s, "Don't read every line. Say: 'We'll start with your world and the problem, move through our solution and "
              "how we build it, then spend real time on the live product — the AI agent and the maps — before "
              "closing on timeline, documentation and deliverables.' Promise the screens are real and you'll demo logic.")

# =========================================================================
# SLIDE 3 — 01 DLD
# =========================================================================
def s_dld():
    s = slide(); bg(s, WHITE)
    header(s, "01 — Context & Mandate", "Dubai Land Department", "01")
    _, tf = tb(s, 0.9, 1.66, 11.5, 0.7)
    para(tf, "DLD registers, regulates and underwrites confidence in one of the world’s most active "
              "real-estate markets. IRETP is the digital backbone for the next decade of that mandate.",
         13.5, INK, font=BODY, first=True, line=1.2)
    pillars = [
        ("Transparency", "A single, trusted source of truth for the public, investors and regulators.", BRIGHT),
        ("Oversight", "Early detection of developer, project and escrow risk before it harms the market.", GREEN),
        ("Investor confidence", "Open data, fair indices and developer accountability — by design.", GREEN2),
        ("Global standing", "Strengthen Dubai’s position on international transparency indices (GRETI).", GOLD2),
    ]
    x0, y0, cw, ch, g = 0.9, 2.62, 2.79, 1.62, 0.13
    for i, (t, d, c) in enumerate(pillars):
        x = x0 + i * (cw + g)
        rect(s, x, y0, cw, ch, fill=CARD, rad=0.08)
        rect(s, x, y0, cw, 0.09, fill=c, rad=0)
        node(s, x + 0.28, y0 + 0.32, 0.14, c)
        _, ct = tb(s, x + 0.28, y0 + 0.56, cw - 0.5, ch - 0.6)
        para(ct, t, 14.5, INK, font=SEMI, first=True)
        para(ct, d, 11, MUTE, font=BODY, before=5, line=1.16)
    # the challenge strip
    rect(s, 0.9, 4.55, 11.53, 1.95, fill=DARK, rad=0.05)
    _, st = tb(s, 1.3, 4.78, 4.4, 1.6, anchor=MSO_ANCHOR.MIDDLE)
    para(st, "THE CHALLENGE IRETP SOLVES", 11, GOLD, font=SEMI, first=True, spacing=2.0)
    para(st, "Fragmented data, manual oversight and limited public transparency — unified into one "
              "intelligent, bilingual platform.", 14, WHITE, font=BODY, before=8, line=1.2)
    frm = [("Fragmented data sources", "One Fabric-aligned data plane"),
           ("Manual risk monitoring", "Automated Early-Warning Risk System"),
           ("Limited public access", "Open data + Arabic-first AI agent")]
    fx = 6.0
    for i, (a, b) in enumerate(frm):
        y = 4.74 + i * 0.535
        _, ft = tb(s, fx, y, 6.2, 0.5, anchor=MSO_ANCHOR.MIDDLE)
        runs(para(ft, "", 11.5, WHITE, first=True)[0], [
            (a, 11.5, PALE, BODY, False), ("    →   ", 11.5, GOLD, SEMI, False), (b, 11.5, WHITE, SEMI, False)])
    footer(s)
    notes(s, "Anchor in DLD's world — show you understand their mandate, not just technology. Four pillars: "
              "transparency, oversight, investor confidence, global standing (GRETI). Then the dark strip frames the "
              "before/after: the three things that are hard today and exactly how IRETP fixes each. Keep it crisp.")

# =========================================================================
# SLIDE 4 — 02 PROJECT SUMMARY
# =========================================================================
def s_summary():
    s = slide(); bg(s, WHITE)
    header(s, "02 — Executive Overview", "Project Summary", "02")
    _, tf = tb(s, 0.9, 1.7, 7.0, 1.0)
    para(tf, "IRETP is a single, bilingual, AI-enabled platform that unifies public transparency and internal "
              "regulatory oversight — and it already exists as a working reference build.",
         15, INK, font=SEMI, first=True, line=1.2)
    rows = [
        ("Two platforms, one system", "An External Public Portal for investors, the public and open-data consumers, "
         "and an Internal Platform for DLD regulatory oversight — cleanly separated, one codebase."),
        ("AI at the core", "An Arabic-first real-estate AI agent, grounded in DLD data, governed, auditable and "
         "UAE-resident — not a bolt-on chatbot."),
        ("Already built", "~90% functional coverage demonstrated today on .NET 9 Clean Architecture, with 106 "
         "automated tests passing — 0 errors, 0 warnings."),
        ("Three months to production", "A de-risked, gated path from this reference build to a UAT-certified, "
         "DESC-aligned, production go-live in the UAE."),
    ]
    y = 2.92
    for i, (t, d) in enumerate(rows):
        rect(s, 0.9, y, 7.0, 0.92, fill=(MIST if i % 2 == 0 else CARD), rad=0.06)
        node(s, 1.16, y + 0.2, 0.13, BRIGHT)
        _, rt = tb(s, 1.5, y + 0.12, 6.25, 0.72, anchor=MSO_ANCHOR.MIDDLE)
        para(rt, t, 13, INK, font=SEMI, first=True)
        para(rt, d, 10.5, MUTE, font=BODY, before=2, line=1.12)
        y += 1.0
    # side panel
    rect(s, 8.15, 1.7, 4.28, 4.92, fill=DARK, rad=0.05)
    _, pt = tb(s, 8.5, 2.0, 3.6, 1.0)
    para(pt, "WHY IT WINS", 11, GOLD, font=SEMI, first=True, spacing=2.4)
    para(pt, "We don’t ask DLD to imagine the platform.", 15.5, WHITE, font=SEMI, before=8, line=1.14)
    stat = [("~90%", "RFP functional coverage in the build today"),
            ("106", "automated tests passing"),
            ("100%", "UAE data residency"),
            ("3 mo", "to a production go-live")]
    yy = 3.45
    for v, d in stat:
        _, vt = tb(s, 8.5, yy, 3.6, 0.78)
        runs(para(vt, "", 26, WHITE, first=True)[0], [(v, 26, BRIGHT, HEAVY, False)])
        para(vt, d, 10.5, PALE, font=BODY, before=0, line=1.05)
        yy += 0.79
    footer(s)
    notes(s, "This is your thesis slide. Land the headline sentence slowly. Walk the four rows left-to-right; spend "
              "most energy on 'Already built' and '106 tests'. The dark panel is your closing argument in miniature: "
              "'We don't ask DLD to imagine the platform.' Point at ~90% and 100% UAE. Everything after this is proof.")

# =========================================================================
# SLIDE 5 — 03 REQUIREMENTS
# =========================================================================
def s_requirements():
    s = slide(); bg(s, WHITE)
    header(s, "03 — Understanding the Requirement", "One platform, four audiences", "03")
    _, tf = tb(s, 0.9, 1.66, 11.5, 0.5)
    para(tf, "The RFP asks one system to serve four very different audiences at the same time. IRETP is "
              "structured around exactly these four.", 13, MUTE, font=BODY, first=True, line=1.15)
    aud = [
        ("Public & Investors", "Price & rental indices · project search · developer scorecards · "
         "GIS heatmaps · Open Data APIs · AI agent (AR/EN)", BRIGHT),
        ("DLD Regulatory Staff", "Early-Warning Risk System · escrow monitoring · AML / beneficial "
         "ownership · immutable audit trails · regulatory reports", GREEN),
        ("Developers & Partners", "Partner APIs · project & registration status · escrow compliance · "
         "benchmark data · self-service developer portal", GREEN2),
        ("Ministry & Government", "KPI dashboards · ESG indicators · macro indices · international "
         "benchmarking · GRETI sub-index tracking", GOLD2),
    ]
    x0, y0, cw, ch, gx, gy = 0.9, 2.45, 5.72, 1.7, 0.1, 0.16
    for i, (t, d, c) in enumerate(aud):
        x, y = x0 + (i % 2) * (cw + gx), y0 + (i // 2) * (ch + gy)
        rect(s, x, y, cw, ch, fill=CARD, rad=0.07)
        rect(s, x, y, 0.08, ch, fill=c)
        _, ct = tb(s, x + 0.36, y + 0.24, cw - 0.6, ch - 0.4)
        para(ct, t, 15, INK, font=SEMI, first=True)
        para(ct, d, 11.5, MUTE, font=BODY, before=8, line=1.28)
    # non-negotiables strip
    _, nt = tb(s, 0.9, 6.42, 2.4, 0.5, anchor=MSO_ANCHOR.MIDDLE)
    para(nt, "NON-NEGOTIABLES", 10.5, GREEN, font=SEMI, first=True, spacing=1.6)
    labels = ["100% UAE residency", "Arabic parity", "DESC security", "Transparency by design"]
    lx = 3.25
    for lab in labels:
        w = 0.3 + len(lab) * 0.092
        chip(s, lx, 6.46, w, 0.42, lab, MIST, GREEN, size=10.5)
        lx += w + 0.16
    footer(s)
    notes(s, "Show you read the RFP as DLD experiences it, not as a feature list. Four audiences, one system. "
              "For each, name two capabilities you'll demo later. Close on the non-negotiables chips — these are "
              "the constraints every decision in the deck respects: UAE residency, Arabic parity, DESC, transparency.")

# =========================================================================
# SLIDE 6 — 04 ARCHITECTURE DIVIDER (hero image)
# =========================================================================
def divider(num, title, subtitle, image, note):
    s = slide()
    pic(s, os.path.join(ASSETS, image), 0, 0, SW, SH)
    full = rect(s, 0, 0, SW, SH); gradient(full, [(0, DARK, 42), (55, DARK, 56), (100, DARK, 84)])
    left = rect(s, 0, 0, 9.7, SH); gradient(left, [(0, DARK, 90), (62, DARK, 55), (100, DARK, 0)], ang=0)
    _, tf = tb(s, 0.95, 2.5, 8.0, 3.0)
    para(tf, "SECTION " + num, 13, GOLD, font=SEMI, first=True, spacing=3.0)
    para(tf, title, 40, WHITE, font=SEMI, before=10, line=1.03)
    para(tf, subtitle, 15.5, PALE, font=BODY, before=12, line=1.2)
    notes(s, note)
    return s

# =========================================================================
# SLIDE 7 — 04 ARCHITECTURE LAYERS
# =========================================================================
def s_arch_layers():
    s = slide(); bg(s, WHITE)
    header(s, "04 — Solution Architecture", "Clean Architecture, three independent hosts", "04")
    # hosts row
    hosts = [("IRETP.Web", "Blazor Server portal — public + admin console"),
             ("IRETP.WebAPI", "Public & investor API · Open Data · SignalR"),
             ("IRETP.AdminAPI", "Internal DLD operations — policy-gated, separate host")]
    x0, y, cw, g = 0.9, 1.78, 3.71, 0.2
    for i, (t, d) in enumerate(hosts):
        x = x0 + i * (cw + g)
        rect(s, x, y, cw, 0.92, fill=GREEN, rad=0.07)
        _, ht = tb(s, x + 0.2, y + 0.12, cw - 0.4, 0.7, anchor=MSO_ANCHOR.MIDDLE)
        para(ht, t, 14.5, WHITE, font=SEMI, first=True)
        para(ht, d, 9.5, PALE, font=BODY, before=2, line=1.05)
    def band(yy, t, d, fill, tcolor, dcolor, h=0.74):
        rect(s, 0.9, yy, 11.53, h, fill=fill, rad=0.05)
        _, bt = tb(s, 1.2, yy, 11.0, h, anchor=MSO_ANCHOR.MIDDLE)
        runs(para(bt, "", 13, tcolor, first=True)[0], [(t + "   ", 13.5, tcolor, SEMI, False), (d, 11, dcolor, BODY, False)])
    band(2.92, "Application", "CQRS use-cases via MediatR — commands, queries, validation, handlers", MIST, INK, MUTE)
    band(3.78, "Domain", "DDD entities, rules & enums — zero framework dependencies (the canonical model)", "E2EFE7", INK, MUTE)
    band(4.64, "Infrastructure", "EF Core 9 · Hangfire · ASP.NET Identity + OIDC · SignalR · AI orchestration", "D7E8DD", INK, MUTE)
    # data row
    rect(s, 0.9, 5.62, 11.53, 0.92, fill=DARK, rad=0.05)
    _, dl = tb(s, 1.2, 5.66, 2.3, 0.84, anchor=MSO_ANCHOR.MIDDLE)
    para(dl, "DATA PLANE", 10.5, GOLD, font=SEMI, first=True, spacing=1.6)
    para(dl, "UAE North · UAE Central DR", 9, PALE, font=BODY, before=2)
    res = ["Azure SQL", "Microsoft Fabric / OneLake (Gold)", "Key Vault (HSM)", "GRS Storage", "App Insights"]
    rx = 3.55
    for r in res:
        w = 0.3 + len(r) * 0.082
        chip(s, rx, 5.86, w, 0.44, r, "0E4630", "FFFFFF", size=9.5)
        rx += w + 0.12
    footer(s)
    notes(s, "Explain the value of Clean Architecture in one breath: the Domain has zero framework dependencies, so "
              "business rules are testable and survive technology change. Three independent hosts mean DLD scales public "
              "traffic separately from internal staff, and internal endpoints are never co-hosted with the public surface. "
              "Bottom line: everything lands in UAE regions, reading DLD's existing Fabric/OneLake Gold layer.")

# =========================================================================
# SLIDE 8 — 04 ARCHITECTURE STACK + FABRIC + DEPLOYMENT
# =========================================================================
def s_arch_stack():
    s = slide(); bg(s, WHITE)
    header(s, "04 — Solution Architecture", "Proven stack, on DLD’s Microsoft Fabric ecosystem", "04")
    # technology stack chips
    _, tf = tb(s, 0.9, 1.66, 11.5, 0.3)
    para(tf, "TECHNOLOGY STACK", 11, GREEN, font=SEMI, first=True, spacing=2.0)
    stack = [".NET 9 LTS", "Blazor Server", "EF Core 9", "MediatR (CQRS)", "Hangfire", "SignalR",
             "ASP.NET Identity + OIDC", "UAE Pass", "Azure", "Bicep IaC", "Docker", "GitHub Actions CI"]
    cx, cy = 0.9, 2.06
    for c in stack:
        w = 0.34 + len(c) * 0.083
        if cx + w > 12.43:
            cx = 0.9; cy += 0.52
        chip(s, cx, cy, w, 0.42, c, CARD, INK, size=10.5)
        rect(s, cx, cy, 0.07, 0.42, fill=BRIGHT)
        cx += w + 0.14
    # three panels
    panels = [
        ("Microsoft Fabric / OneLake", "Reads DLD’s existing Gold layer through a single abstraction "
         "(IFabricGoldDataSource) — no parallel ETL, no parallel data store. (RFP v1.3 §11.4)"),
        ("Deployment", "Azure UAE North (primary) + UAE Central (warm DR) · Front Door + WAF · private "
         "endpoints · P1v3 multi-instance · RPO ≤ 15 min, RTO ≤ 1 h."),
        ("Visualisation", "Power BI for the Internal Platform; a custom-web layer (Chart.js + MapLibre) for the "
         "External Portal — AED 0 licensing at public scale. (RFP v1.3 §11.4.2)"),
    ]
    x0, y0, cw, ch, g = 0.9, 3.78, 3.71, 2.55, 0.2
    cols = [BRIGHT, GREEN, GOLD2]
    for i, (t, d) in enumerate(panels):
        x = x0 + i * (cw + g)
        rect(s, x, y0, cw, ch, fill=CARD, rad=0.06)
        rect(s, x, y0, cw, 0.1, fill=cols[i])
        _, pt = tb(s, x + 0.26, y0 + 0.32, cw - 0.5, ch - 0.5)
        para(pt, t, 14, INK, font=SEMI, first=True, line=1.05)
        para(pt, d, 11, MUTE, font=BODY, before=8, line=1.32)
    footer(s)
    notes(s, "Two messages. First, the stack is proven and current (.NET 9 LTS) — no exotic bets. Second, and most "
              "important for v1.3: we align to DLD's existing Microsoft Fabric / OneLake. We read your Gold layer through "
              "one abstraction — no parallel ETL. Mention the Power BI vs custom-web analysis: Power BI internally, "
              "custom-web externally to avoid per-user licensing at public scale. All UAE, with warm DR.")

# =========================================================================
# SLIDE 9 — 05 METHODOLOGY
# =========================================================================
def s_methodology():
    s = slide(); bg(s, WHITE)
    header(s, "05 — Delivery Methodology", "How we build — and how we assure", "05")
    items = [
        ("Engineering discipline", "Clean Architecture + DDD + CQRS. A framework-agnostic core that is testable, "
         "maintainable and survives technology change.", BRIGHT),
        ("Agile, with hard gates", "Three monthly milestones with G1 / G2 / G3 acceptance gates and fortnightly DLD "
         "steering — progress you can audit.", GREEN),
        ("Quality built in", "106 automated tests today; unit, integration, E2E, performance, security and AI-accuracy "
         "checks run in CI on every change.", GREEN2),
        ("AI governance", "Tiered model routing, UAE-residency metadata on every call, deterministic analytics overlay, "
         "an advisory guardrail and an EN/AR accuracy harness.", GOLD2),
        ("Security by design", "DESC ISR v3 control mapping, OWASP, external VAPT, 24×7 SOC and an immutable audit "
         "trail — engineered from day one, not retrofitted.", GREEN),
        ("Data familiarisation", "A 2-week discovery sprint maps every DLD source system to the Gold-layer schema "
         "before a single number is published. (§12.1)", BRIGHT),
    ]
    x0, y0, cw, ch, gx, gy = 0.9, 1.84, 3.71, 1.46, 0.2, 0.2
    for i, (t, d, c) in enumerate(items):
        x, y = x0 + (i % 3) * (cw + gx), y0 + (i // 3) * (ch + gy)
        rect(s, x, y, cw, ch, fill=CARD, rad=0.07)
        node(s, x + 0.26, y + 0.28, 0.14, c)
        _, ct = tb(s, x + 0.26, y + 0.52, cw - 0.5, ch - 0.6)
        para(ct, t, 13.5, INK, font=SEMI, first=True)
        para(ct, d, 10.5, MUTE, font=BODY, before=5, line=1.2)
    rect(s, 0.9, 6.5, 11.53, 0.0)  # spacer noop
    footer(s)
    notes(s, "Methodology is where you separate from competitors who'll show a generic Agile slide. Hit three "
              "differentiators hard: (1) the architecture makes quality cheap — 106 tests already; (2) AI is governed, "
              "not vibes — residency metadata, guardrail, accuracy harness; (3) security and data familiarisation are "
              "engineered in. The G1/G2/G3 gates give DLD audit points to release payment against.")

# =========================================================================
# SLIDE 10 — 06 FUNCTIONAL REQUIREMENTS
# =========================================================================
def s_functional():
    s = slide(); bg(s, WHITE)
    header(s, "06 — Functional Coverage", "What the platform does — today", "06")
    _, tf = tb(s, 0.9, 1.66, 11.5, 0.4)
    runs(para(tf, "", 12.5, MUTE, first=True)[0], [
        ("Every group below is ", 12.5, MUTE, BODY, False),
        ("implemented in the reference build", 12.5, GREEN, SEMI, False),
        (" — not roadmap, not mockup.", 12.5, MUTE, BODY, False)])
    groups = [
        ("External Public Portal", "CMS · KPI dashboard · transactions · GIS map · price & rental index · slice-and-dice analytics"),
        ("Real Estate AI Agent", "RAG-grounded · multi-model orchestration · bilingual · guardrails · cross-session memory"),
        ("Investor Tools", "Watchlist · alerts (email / SMS / in-app) · saved views · personalised scorecard PDF"),
        ("Early-Warning Risk System", "10 risk indicators · L1–L4 auto-escalation · SLA tracking · risk heatmap · playbooks"),
        ("Developer Rating & Escrow", "8-criteria scoring · configurable weights · escrow monitoring · monthly health reports"),
        ("Open Data Portal", "OpenAPI 3.0 contract · tiered rate limits · self-service developer portal & API keys"),
        ("Transparency Plus", "Beneficial ownership · mortgage analytics · ESG indicators · GRETI sub-index tracker"),
        ("Multilingual & RTL", "Arabic-first EN/AR parity in Phase 1 · six further languages · full RTL mirroring"),
    ]
    x0, y0, cw, ch, gx, gy = 0.9, 2.28, 5.72, 0.96, 0.1, 0.14
    for i, (t, d) in enumerate(groups):
        x, y = x0 + (i % 2) * (cw + gx), y0 + (i // 2) * (ch + gy)
        rect(s, x, y, cw, ch, fill=CARD, rad=0.07)
        chip(s, x + cw - 1.38, y + 0.16, 1.18, 0.34, "IMPLEMENTED", MIST, GREEN, size=8)
        rect(s, x + cw - 1.38, y + 0.16, 0.06, 0.34, fill=BRIGHT)
        _, ct = tb(s, x + 0.28, y + 0.12, cw - 1.7, ch - 0.2, anchor=MSO_ANCHOR.MIDDLE)
        para(ct, t, 13, INK, font=SEMI, first=True)
        para(ct, d, 9.8, MUTE, font=BODY, before=2, line=1.12)
    footer(s)
    notes(s, "Don't read the matrix — let it land visually. Say: 'Eight capability groups the RFP asks for; every "
              "one already carries an IMPLEMENTED badge.' Then pivot: 'Rather than talk through them, let me show you the "
              "real screens.' This sets up the showcase — the strongest part of the meeting.")

# =========================================================================
# SLIDE 12+ — SCREENS
# =========================================================================
def s_screen(key, eyebrow, title, points, note, badge=None):
    s = slide(); bg(s, WHITE)
    header(s, eyebrow, title, "07")
    img_w = 7.55
    img_h = img_w / SHOT_AR
    iy = 1.95
    pic(s, os.path.join(FRAMED, key + ".png"), 0.9, iy, img_w, img_h, shadowed=True)
    px = 0.9 + img_w + 0.45
    pw = SW - px - 0.9
    _, pt = tb(s, px, 2.0, pw, 4.4)
    if badge:
        para(pt, badge.upper(), 10, GOLD2, font=SEMI, first=True, spacing=1.8)
    first = badge is None
    for i, (h, d) in enumerate(points):
        p, _ = para(pt, h, 13, INK, font=SEMI, first=(first and i == 0), before=(0 if (first and i == 0) else 12))
        para(pt, d, 10.8, MUTE, font=BODY, before=3, line=1.22)
    footer(s)
    notes(s, note)

# =========================================================================
# SLIDE — SCREENS MONTAGE
# =========================================================================
def s_montage():
    s = slide(); bg(s, WHITE)
    header(s, "07 — The Platform", "… and the rest of the transparency suite", "07")
    _, tf = tb(s, 0.9, 1.66, 11.5, 0.4)
    para(tf, "Beyond the four headline screens, the full suite is built and running.", 12.5, MUTE, font=BODY, first=True)
    grid = [("transactions", "Transactions Explorer", "Filter · sort · export (CSV / PDF)"),
            ("scorecards", "Developer Scorecards", "Public, weighted, transparent"),
            ("bo", "Beneficial Ownership", "AML disclosure · access-audited"),
            ("mortgage", "Mortgage Analytics", "LTV, tenor & rate transparency"),
            ("audit", "Immutable Audit Logs", "Append-only · every action")]
    cw = 2.5
    ch = cw / SHOT_AR
    r0y, r1y = 2.46, 2.46 + ch + 0.5 + 0.3
    positions = [(2.665, r0y), (5.415, r0y), (8.165, r0y), (4.04, r1y), (6.79, r1y)]
    for (x, y), (key, t, d) in zip(positions, grid):
        pic(s, os.path.join(FRAMED, key + ".png"), x, y, cw, ch, shadowed=True)
        _, ct = tb(s, x + 0.02, y + ch + 0.07, cw - 0.04, 0.5)
        para(ct, t, 11.5, INK, font=SEMI, first=True)
        para(ct, d, 9.5, MUTE, font=BODY, before=1, line=1.0)
    footer(s)
    notes(s, "Quick montage — don't dwell. 'Beyond the four headline screens, the full transparency suite is built: "
              "transactions, public developer scorecards, beneficial ownership with access auditing, mortgage analytics, "
              "and an append-only audit log.' One sentence each, then move to non-functionals.")

# =========================================================================
# SLIDE — 08 NON-FUNCTIONAL
# =========================================================================
def s_nfr():
    s = slide(); bg(s, WHITE)
    header(s, "08 — Non-Functional Requirements", "Performance, security and trust", "08")
    quads = [
        ("Performance", BRIGHT, [
            "P95 public page load ≤ 2.0 s", "AI agent P95 response < 3.5 s",
            "KPI dashboard freshness ≤ 15 min", "99.9% public / 99.95% internal availability"]),
        ("Security", GREEN, [
            "MFA + RBAC policy registry", "30-min internal JWTs, refresh rotation",
            "TLS 1.2+ · partner API-key tiers", "Immutable, append-only audit trail"]),
        ("Data residency", GOLD2, [
            "100% UAE — UAE North + UAE Central DR", "AI inference UAE-resident, or flagged",
            "Backups, logs & analytics in-region", "No cross-border processing without waiver"]),
        ("Compliance", GREEN2, [
            "DESC ISR v3 control mapping", "UAE PDPL · OWASP ASVS L2",
            "ISO 27001 roadmap · external VAPT", "24×7 SOC from Month 1"]),
    ]
    x0, y0, cw, ch, gx, gy = 0.9, 1.84, 5.72, 2.28, 0.1, 0.18
    for i, (t, c, rows) in enumerate(quads):
        x, y = x0 + (i % 2) * (cw + gx), y0 + (i // 2) * (ch + gy)
        rect(s, x, y, cw, ch, fill=CARD, rad=0.06)
        rect(s, x, y, 0.09, ch, fill=c)
        _, ht = tb(s, x + 0.34, y + 0.2, cw - 0.5, 0.4)
        para(ht, t, 14.5, INK, font=SEMI, first=True)
        _, lt = tb(s, x + 0.34, y + 0.72, cw - 0.6, ch - 0.85)
        for j, r in enumerate(rows):
            runs(para(lt, "", 11, INK, first=(j == 0), before=(0 if j == 0 else 6))[0], [
                ("▪  ", 11, c, BODY, False), (r, 11, INK, BODY, False)])
    footer(s)
    notes(s, "These are commitments, not aspirations — the SLA endpoint /healthz/sla already measures most of them. "
              "Lead with the two DLD cares about most: 100% UAE residency (point at the third quadrant) and the security "
              "posture (MFA, RBAC, immutable audit). Mention DESC ISR v3 is already mapped clause-by-clause in the docs.")

# =========================================================================
# SLIDE — 09 TIMELINE
# =========================================================================
def s_timeline():
    s = slide(); bg(s, WHITE)
    header(s, "09 — Implementation Timeline", "Three months to a production go-live", "09")
    # week ruler
    rx, ry, rw = 0.9, 1.86, 11.53
    _, wt = tb(s, rx, ry, 2.0, 0.3)
    para(wt, "WEEK", 9, MUTE, font=SEMI, first=True, spacing=1.5)
    for wk in range(1, 13):
        x = rx + 1.4 + (wk - 1) * ((rw - 1.4) / 12.0)
        _, t = tb(s, x, ry, 0.6, 0.3)
        para(t, f"W{wk}", 9, MUTE, font=BODY, first=True)
    line_h(s, rx, ry + 0.34, rw, color=LINE, weight=1.0)
    months = [
        ("MONTH 1", "Foundation & Integration", BRIGHT,
         ["Stand up UAE Azure environments (Bicep)", "UAE Pass federation · integration design",
          "SOC + DESC-CSP engaged · residual UI polish"], "G1", "Prod env ready · contracts signed"),
        ("MONTH 2", "Build, Integrate & Harden", GREEN,
         ["Integrate DLD sources · EWRS on live data", "External VAPT · load & DR testing",
          "AI accuracy harness → ≥ 90% on DLD catalogue"], "G2", "Feature-complete · security-cleared"),
        ("MONTH 3", "UAT, Certify & Go-Live", GOLD2,
         ["UAT waves 1–2 · defect triage", "DESC-CSP audit · VAPT re-test · training",
          "Go/No-Go → Go-Live + 14-day hyper-care"], "G3", "UAT signed · certified · live"),
    ]
    x0, y0, cw, ch, g = 0.9, 2.5, 3.71, 3.55, 0.2
    for i, (m, theme, c, acts, gate, gd) in enumerate(months):
        x = x0 + i * (cw + g)
        rect(s, x, y0, cw, ch, fill=CARD, rad=0.06)
        rect(s, x, y0, cw, 0.78, fill=c, rad=0.06)
        rect(s, x, y0 + 0.5, cw, 0.28, fill=c)  # square off bottom of header
        _, mt = tb(s, x + 0.28, y0 + 0.1, cw - 0.5, 0.62, anchor=MSO_ANCHOR.MIDDLE)
        para(mt, m, 11, WHITE, font=SEMI, first=True, spacing=1.5)
        para(mt, theme, 14, WHITE, font=SEMI, before=1)
        _, at = tb(s, x + 0.3, y0 + 1.0, cw - 0.56, 1.75)
        for j, a in enumerate(acts):
            runs(para(at, "", 11, INK, first=(j == 0), before=(0 if j == 0 else 9))[0], [
                ("•  ", 11, c, SEMI, False), (a, 11, INK, BODY, False, )][:2] if False else [
                ("•  ", 11, c, SEMI, False), (a, 11, INK, BODY, False)])
        # gate badge
        rect(s, x + 0.3, y0 + ch - 0.74, cw - 0.6, 0.52, fill=DARK, rad=0.1)
        _, gt = tb(s, x + 0.5, y0 + ch - 0.72, cw - 0.9, 0.48, anchor=MSO_ANCHOR.MIDDLE)
        runs(para(gt, "", 11, WHITE, first=True)[0], [
            ("GATE " + gate + "   ", 11, GOLD, HEAVY, False), (gd, 9.5, PALE, BODY, False)])
    footer(s)
    notes(s, "The timeline is your risk story. Because the platform is already built, Month 1 isn't 'start coding' — "
              "it's environments, integration design and UAE Pass. Months 2–3 are integration, hardening, VAPT and UAT. "
              "Each month ends in a gate (G1/G2/G3) that DLD signs before the next payment. Emphasise: 12 weeks is "
              "credible precisely because we're not starting from zero.")

# =========================================================================
# SLIDE — 10 DOCUMENTATION
# =========================================================================
def s_docs():
    s = slide(); bg(s, WHITE)
    header(s, "10 — Documentation Provided", "A complete, living knowledge pack", "10")
    _, tf = tb(s, 0.9, 1.66, 11.5, 0.4)
    para(tf, "Every document below already exists in the reference repository — delivered, version-controlled, "
              "and self-validating: CI tests prove the docs still match the code.", 12.5, MUTE, font=BODY, first=True, line=1.18)
    docs = [
        "Architecture & Integration Map", "Compliance Matrix (RFP → code → test)", "DESC ISR v3 Control Mapping",
        "OWASP Top 10 Mapping", "STRIDE Threat Model", "Risk Register",
        "Security & Compliance Plan", "UAT Plan (10 categories)", "Operations Runbook",
        "Disaster Recovery (RTO/RPO)", "Data Dictionary", "API Reference + OpenAPI 3.0",
        "Data Familiarisation Pack", "Knowledge Transfer & Exit Plan", "Visualisation Layer Analysis",
    ]
    x0, y0, cw, ch, gx, gy = 0.9, 2.5, 3.71, 0.74, 0.2, 0.16
    for i, d in enumerate(docs):
        x, y = x0 + (i % 3) * (cw + gx), y0 + (i // 3) * (ch + gy)
        rect(s, x, y, cw, ch, fill=CARD, rad=0.08)
        rect(s, x, y, 0.07, ch, fill=BRIGHT)
        _, dt = tb(s, x + 0.3, y, cw - 0.45, ch, anchor=MSO_ANCHOR.MIDDLE)
        para(dt, d, 11.5, INK, font=SEMI, first=True, line=1.05)
    footer(s)
    notes(s, "Most bidders promise documentation later. You hand over fifteen documents that already exist — and the "
              "killer line: they're self-validating. A CI test fails if a doc references a file that no longer exists, so "
              "the documentation can't silently rot. Call out the Compliance Matrix: every RFP requirement mapped to the "
              "implementing file and its test.")

# =========================================================================
# SLIDE — 11 DELIVERABLES
# =========================================================================
def s_deliverables():
    s = slide(); bg(s, WHITE)
    header(s, "11 — Project Deliverables", "What the Dubai Land Department receives", "11")
    items = [
        ("The Platform", "External Portal + Internal Platform + Open Data API — bilingual, production-deployed in the UAE.", BRIGHT),
        ("Source & IP", "Full source code, Bicep IaC and CI/CD pipelines — all intellectual property vests in DLD.", GREEN),
        ("Documentation", "The complete knowledge pack — architecture, compliance, security, ops, DR and data.", GREEN2),
        ("Assurance", "UAT sign-off, VAPT report + remediation, DESC-CSP certification, performance & DR test reports.", GOLD2),
        ("Enablement", "Super-user training, ops runbook walkthrough, developer portal, and 14-day hyper-care.", GREEN),
        ("Support", "12-month warranty, SLA-backed (Sev 1: 1 h response / 4 h restore), 24×7 SOC.", BRIGHT),
    ]
    x0, y0, cw, ch, gx, gy = 0.9, 1.84, 5.72, 1.3, 0.1, 0.16
    for i, (t, d, c) in enumerate(items):
        x, y = x0 + (i % 2) * (cw + gx), y0 + (i // 2) * (ch + gy)
        rect(s, x, y, cw, ch, fill=CARD, rad=0.07)
        rect(s, x, y, 0.08, ch, fill=c)
        _, ct = tb(s, x + 0.34, y + 0.18, cw - 0.55, ch - 0.3, anchor=MSO_ANCHOR.MIDDLE)
        para(ct, t, 13.5, INK, font=SEMI, first=True)
        para(ct, d, 11, MUTE, font=BODY, before=4, line=1.2)
    # payment/milestone strip
    rect(s, 0.9, 6.18, 11.53, 0.66, fill=DARK, rad=0.06)
    _, mt = tb(s, 1.15, 6.18, 2.2, 0.66, anchor=MSO_ANCHOR.MIDDLE)
    para(mt, "MILESTONES", 10, GOLD, font=SEMI, first=True, spacing=1.5)
    miles = ["Kickoff", "Gate G1", "Gate G2", "UAT sign-off", "Go-Live", "Hyper-care"]
    mx = 3.3
    for i, m in enumerate(miles):
        _, t = tb(s, mx, 6.18, 1.7, 0.66, anchor=MSO_ANCHOR.MIDDLE)
        runs(para(t, "", 10.5, WHITE, first=True)[0], [(m, 10.5, WHITE, SEMI, False)])
        if i < len(miles) - 1:
            _, a = tb(s, mx + 1.36, 6.18, 0.3, 0.66, anchor=MSO_ANCHOR.MIDDLE)
            para(a, "→", 12, GOLD, font=BODY, first=True)
        mx += 1.55
    footer(s)
    notes(s, "Make ownership explicit: all IP vests in DLD — source, IaC, pipelines, derived data. Then the assurance "
              "line is what a government buyer needs to defend the decision internally: VAPT, DESC-CSP certification, UAT "
              "sign-off. Close on support: 12-month warranty, 1-hour Sev-1 response, 24×7 SOC. The strip shows the "
              "milestone path payments attach to.")

# =========================================================================
# SLIDE — CLOSING
# =========================================================================
def s_closing():
    s = slide()
    pic(s, os.path.join(ASSETS, "closing_hero.png"), 0, 0, SW, SH)
    full = rect(s, 0, 0, SW, SH); gradient(full, [(0, DARK, 48), (55, DARK, 60), (100, DARK, 86)])
    left = rect(s, 0, 0, 9.9, SH); gradient(left, [(0, DARK, 90), (62, DARK, 60), (100, DARK, 0)], ang=0)
    _, tf = tb(s, 0.95, 1.5, 9.8, 1.2)
    para(tf, "IN CLOSING", 13, GOLD, font=SEMI, first=True, spacing=3.0)
    para(tf, "Why [Your Company]", 38, WHITE, font=SEMI, before=8)
    pts = [
        ("Already built", "We’re not promising a platform — we’re demonstrating one. ~90% RFP coverage, today."),
        ("UAE-native & compliant", "100% data residency, DESC ISR v3 mapped, PDPL-ready, Arabic-first by construction."),
        ("De-risked delivery", "A 3-month, gated path to go-live, backed by 106 passing tests and a complete doc pack."),
    ]
    y = 3.2
    for t, d in pts:
        node(s, 0.98, y + 0.08, 0.16, BRIGHT)
        _, ct = tb(s, 1.32, y, 8.0, 0.9)
        runs(para(ct, "", 16, WHITE, first=True)[0], [(t + "   ", 16, WHITE, SEMI, False), (d, 13, PALE, BODY, False)])
        y += 0.86
    line_h(s, 0.98, 6.35, 5.6, color=GOLD, weight=1.4)
    _, et = tb(s, 0.95, 6.5, 11.0, 0.8)
    para(et, "Thank you — we would welcome your questions.", 17, WHITE, font=SEMI, first=True)
    para(et, "[Your Company Name]   ·   [Presenter, Title]   ·   [email]   ·   [phone]", 11, PALE, font=BODY, before=4)
    notes(s, "Bring it home in three beats: already built, UAE-native & compliant, de-risked. Then stop talking and "
              "invite questions warmly. If they engage, offer a live walkthrough of the AI agent or the EWRS escalation "
              "ladder on the running system — that interactivity is what seals it. Replace the bidder/contact "
              "placeholders before the meeting.")

# ---- Build ---------------------------------------------------------------
s_cover()
s_agenda()
s_dld()
s_summary()
s_requirements()
divider("04", "Solution Architecture", "Clean Architecture, three independent hosts — aligned to DLD’s "
        "existing Microsoft Fabric / OneLake ecosystem.", "arch_hero.png",
        "Transition slide. One line: 'Here's how it's engineered.' Let the visual breathe for a beat, then move to the "
        "layer diagram. Keep it to ~10 seconds.")
s_arch_layers()
s_arch_stack()
s_methodology()
s_functional()
divider("07", "The Platform — Live Screens", "Not mockups. A running product — the AI agent, the maps, and the "
        "full transparency suite.", "aimaps_hero.png",
        "The pivot to proof. Say it plainly: 'Everything you're about to see is a running system, not a design comp.' "
        "This is the emotional peak of the meeting — slow down and let them look.")
s_screen("dashboard", "07 — The Platform", "Market Pulse — the public dashboard",
         [("Live KPIs, 15-minute freshness", "Total volume, transactions, average PSF and top zone — sourced from the Fabric Gold layer."),
          ("FR-004 market charts", "Monthly volume, off-plan vs ready mix, top-5 zones and a price-per-sqft trend with crosshair hover."),
          ("Bilingual, branded, fast", "Full EN/AR with RTL mirroring; P95 page load under 2 seconds.")],
         "The dashboard is the public's front door. Point out the freshness badge ('as of 4 min ago') — that's the "
         "15-minute SLA, live. The four charts are FR-004, fully interactive. Note it's the same DLD-green design system "
         "as the real product.", badge="The public's front door")
s_screen("map", "07 — The Platform  ·  MAPS", "Interactive GIS Atlas",
         [("FR-009 — four heatmap modes", "Volume, price-per-sqft, yield and ESG coverage — banded on the official Dubai zone GeoJSON."),
          ("FR-011 — individual project pins", "Every active project as a pin, colour-coded by status: completed, under construction, future, stalled."),
          ("Click-through detail", "Select a zone for average PSF, unit count and yield; click a pin for project-level detail.")],
         "This is the 'Maps' your headline calls out. Switch mentally between the four heatmap modes and the project "
         "pins. The data sits on the OFFICIAL Dubai zone GeoJSON — not a cartoon map. 89 zones, 500+ projects in the "
         "seed. Offer to drive it live in Q&A.", badge="Headline 7 · Maps")
s_screen("ai", "07 — The Platform  ·  AI", "Real Estate AI Agent",
         [("Arabic-first, RAG-grounded", "Answers grounded in DLD data with source citations — every reply is traceable, not generative guesswork."),
          ("Deterministic analytics overlay", "Trend, anomaly and correlation are computed in-process and quoted, so the model reports numbers it doesn’t invent."),
          ("Governed by design", "Non-advisory guardrail refuses investment advice; UAE-resident inference; multi-model routing with fallback.")],
         "The centrepiece. Three things make this enterprise-grade, not a chatbot: (1) it cites its sources; (2) the "
         "numbers come from deterministic analytics, not the model's imagination; (3) a guardrail refuses regulated "
         "investment advice in AR and EN. Read the Arabic example aloud if comfortable. Mention UAE-resident inference.",
         badge="Headline 7 · AI")
s_screen("ewrs", "07 — The Platform  ·  INTERNAL", "Early-Warning Risk System",
         [("10-indicator risk engine", "Delivery delay, escrow shortfall, construction suspension, volume decline, developer deterioration and more."),
          ("L1–L4 auto-escalation", "From Project Officer to DG + Regulator, on UAE business-hours SLAs, with a live escalation ladder."),
          ("Immutable audit, by law", "Every threshold change and acknowledgement is written to an append-only audit log (§10.2).")],
         "Switch register here — this is for DLD's regulators, not the public. The risk engine runs every 15 minutes "
         "and auto-escalates through four levels on business-hours SLAs. The audit log is append-only at the database "
         "boundary — it physically cannot be edited. This is the oversight backbone.", badge="Internal oversight")
s_montage()
s_nfr()
s_timeline()
s_docs()
s_deliverables()
s_closing()

prs.save(OUT)
print("Saved:", OUT, "slides:", len(prs.slides._sldIdLst))
